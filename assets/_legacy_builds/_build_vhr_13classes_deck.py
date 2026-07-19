#!/usr/bin/env python3
"""Per-class deck: detecting Thomas's 13 waste/material classes from 30-50 cm VHR satellite (+ drone).

Thomas / internal register: absolute B/N, no colour, solid bullets, English,
self-explanatory titles, refs as title + year. One slide per class.
Source: docs/02_research/sota_vhr_13classes.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x66, 0x66, 0x66)
FONT = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
M = Inches(0.7)


def new():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title(slide, text, sub=None):
    tf = _tb(slide, M, Inches(0.45), EMU_W - 2 * M, Inches(1.0))
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = BLACK
    ln = slide.shapes.add_connector(2, M, Inches(1.32), EMU_W - M, Inches(1.32))
    ln.line.color.rgb = BLACK; ln.line.width = Pt(1.5)
    if sub:
        tf2 = _tb(slide, M, Inches(1.4), EMU_W - 2 * M, Inches(0.5))
        r2 = tf2.paragraphs[0].add_run(); r2.text = sub
        r2.font.name = FONT; r2.font.size = Pt(14); r2.font.italic = True; r2.font.color.rgb = GREY


def bullets(slide, items, top=Inches(1.7), size=18, gap=10):
    tf = _tb(slide, M, top, EMU_W - 2 * M, EMU_H - top - Inches(0.5))
    first = True
    for it in items:
        lvl, txt, bold = 0, it, False
        if isinstance(it, tuple):
            if len(it) == 3:
                lvl, txt, bold = it
            else:
                lvl, txt = it
        if not isinstance(lvl, int):
            lvl = 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        mark = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark + txt
        r.font.name = FONT; r.font.size = Pt(size - 2 * lvl)
        r.font.bold = bold
        r.font.color.rgb = BLACK if lvl == 0 else RGBColor(0x33, 0x33, 0x33)


def footer(slide, text):
    tf = _tb(slide, M, EMU_H - Inches(0.5), EMU_W - 2 * M, Inches(0.4))
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


def table(slide, headers, rows, col_w, top, fsize=12):
    g = slide.shapes.add_table(len(rows) + 1, len(headers), M, top,
                               Emu(int(sum(col_w))), Inches(0.4 * (len(rows) + 1)))
    t = g.table
    for ci, w in enumerate(col_w):
        t.columns[ci].width = Emu(int(w))
    for ci, h in enumerate(headers):
        c = t.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        c.margin_left = Pt(4); c.margin_right = Pt(4); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        r = c.text_frame.paragraphs[0].add_run(); r.text = h
        r.font.name = FONT; r.font.size = Pt(fsize); r.font.bold = True; r.font.color.rgb = BLACK
    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.margin_left = Pt(4); c.margin_right = Pt(4); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            r = c.text_frame.paragraphs[0].add_run(); r.text = val
            r.font.name = FONT; r.font.size = Pt(fsize); r.font.color.rgb = BLACK
            if ci == 0:
                r.font.bold = True


# ---- Title ----
s = new()
tf = _tb(s, M, Inches(2.3), EMU_W - 2 * M, Inches(2.2))
r = tf.paragraphs[0].add_run()
r.text = "Detecting waste materials from 30–50 cm satellite imagery"
r.font.name = FONT; r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = BLACK
p = tf.add_paragraph(); p.space_before = Pt(14)
r = p.add_run(); r.text = "What the literature does for our 13 target classes — by satellite and by drone"
r.font.name = FONT; r.font.size = Pt(19); r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
ln = s.shapes.add_connector(2, M, Inches(4.4), Inches(6.5), Inches(4.4))
ln.line.color.rgb = BLACK; ln.line.width = Pt(1.5)
footer(s, "PERIVALLON thesis · input = VHR 30–50 cm satellite, resolution-first · 2026-06")

# ---- Framing ----
s = new()
title(s, "The rule that decides feasibility: object vs material")
bullets(s, [
    "At 30–50 cm broadband VNIR (no SWIR), a class is feasible if its identity is in SHAPE / CONTEXT, weak if it is in SPECTRAL chemistry.",
    ("", "Object / shape — detectable: vehicles, tanks, containers, rubble heaps, bulky items, firewood, tires, big bags.", True),
    ("", "Material / spectrum — weak at broadband VHR: asbestos, plastic polymer type, foundry slag, sludge, scrap composition.", True),
    "Confirmed independently by the Waste Management RS survey (2024) and the Gibellini (2025) review.",
], size=18)

# ---- Anchors ----
s = new()
title(s, "Who already does this from satellite", "All detect aggregate dumpsites / binary waste by shape — none does per-class material")
rows = [
    ["Gibellini 2025", "VHR <50 cm (WV-3, GE, AGEA)", "Swin-T + RSP, binary", "F1 92.0% / Acc 94.6%"],
    ["CascadeDumpNet 2024", "Pléiades 0.5 m", "CNN detection + AutoML", "84.6% mAP, transferable"],
    ["Sun 2023", "Gaofen / SuperView 0.3–1 m", "CNN + channel attention", "~1000 dumpsites, 15 countries"],
    ["AerialWaste 2023", "WV-3 0.3 m + GE 0.5 m", "ResNet-50+FPN, binary", "AP 88.0% / 94.5% (0.2 m)"],
    ["Disaitek + Airbus 2024", "Pléiades Neo 30 cm", "AI semantic segmentation", "waste ≥2 m² at ~95% (operational)"],
]
table(s, ["Work (year)", "Platform / GSD", "Method", "Result"],
      rows, [Inches(2.9), Inches(3.4), Inches(2.9), Inches(2.7)], Inches(1.7), fsize=12)
footer(s, "Disaitek/Airbus = vendor claim; it qualifies type: end-of-life vehicles, construction waste, tires, vegetal waste")

# ---- Per-class slides ----
CLASSES = [
    ("1 · Rubble / C&D debris", "object (shape)",
     [("", "Satellite: YES — CWLD (GF-2 80 cm + Google Earth 50 cm), improved DeepLabV3+ segmentation, F1 88.9% / IoU 82% (Beijing, 2024); Yong DeepLabv3+ 1 m, F1 77.4%; Disaitek 30 cm."),
      ("", "Drone: YES — FCN + structure-from-motion, IoU 0.9 for concrete + volume (Drones, 2022)."),
      ("", "Method: semantic segmentation / object detection."),
      ("", "Verdict: strong at 30–50 cm — recognised as heap morphology.", True)],
     "CWLD 2024 · Yong (in Gibellini 2025) · Cheng et al. Drones 2022"),

    ("2 · Foundry waste / slag", "material (spectrum)",
     [("", "Satellite: NO dedicated detector — only AerialWaste annotations (9 labels, rarest)."),
      ("", "Drone: NO — composition studied only at close-range hyperspectral (EJRS 2015)."),
      ("", "Method: none at VHR; spectral phase-mapping in lab."),
      ("", "Verdict: GAP — spectral problem + extreme data scarcity. Weakest class.", True)],
     "Torres & Fraternali 2023 · Iron/steel by-products, Eur. J. Remote Sensing 2015"),

    ("3 · Vehicles / end-of-life vehicles", "object (shape)",
     [("", "Satellite: YES — region-based detector + domain adaptation (+10%, RS 2020); heavy-duty truck classification (2025); Disaitek qualifies end-of-life vehicles at 30 cm."),
      ("", "Drone: YES — standard UAV vehicle detection."),
      ("", "Method: object detection (region-based / YOLO-family)."),
      ("", "Verdict: excellent at 30–50 cm. Generic vehicles mature; a scrapyard-ELV-specific detector is a sub-gap.", True)],
     "Koga et al. RS 2020 · heavy-duty truck PMC 2025 · Disaitek/Airbus 2024"),

    ("4 · Scrap metal", "object + material",
     [("", "Satellite: scrapyard SCENES detectable (AerialWaste, 167 labels) — metal type NOT."),
      ("", "Drone: NO remote detector — composition = close-range infrared (Hybrid-YOLOv5 ELV non-ferrous, 84.2% mAP, not remote sensing)."),
      ("", "Method: scene classification; composition needs spectral/close-range."),
      ("", "Verdict: scenes yes; metal composition is a GAP.", True)],
     "Torres & Fraternali 2023 · Hybrid-YOLOv5, Scientific Reports 2025"),

    ("5 · Bulky items (ingombranti)", "object (shape)",
     [("", "Satellite: YES — AerialWaste, 286 labels (2nd most frequent)."),
      ("", "Drone: partial — large-area UAV solid-waste segmentation (~450 km², >94% OA), but as a generic 'waste pile' class."),
      ("", "Method: scene classification / semantic segmentation."),
      ("", "Verdict: good at 30–50 cm — object; no per-type breakdown.", True)],
     "Torres & Fraternali 2023 · Liu et al. Applied Sciences 14:2084, 2024"),

    ("6 · Containers / skips", "object (shape)",
     [("", "Satellite: YES — CenterNet + Mask R-CNN ensemble classifies container size/status; Satellogic ~90% counting at 0.7 m; AerialWaste 167 labels."),
      ("", "Drone: not specifically needed."),
      ("", "Method: object detection / instance segmentation."),
      ("", "Verdict: strong at 30–50 cm — object.", True)],
     "Heavy-duty/container detection PMC · Satellogic 2023 · Torres & Fraternali 2023"),

    ("7 · Sludge (fanghi)", "material / context",
     [("", "Satellite: NO dedicated detector — only tailings-pond analogues (SSD on Gaofen-1, 90.2% acc); AerialWaste 19 labels."),
      ("", "Drone: NO."),
      ("", "Method: context/colour; composition is spectral."),
      ("", "Verdict: GAP — lagoons visible by context, composition not classifiable at broadband VHR.", True)],
     "Yang et al. RS 12:2626, 2020 (tailings analogue) · Torres & Fraternali 2023"),

    ("8 · Wood and firewood", "object (shape)",
     [("", "Satellite: YES — AerialWaste, 173 labels (4th most frequent)."),
      ("", "Drone: partial — woody-debris volume estimation is airborne."),
      ("", "Method: scene classification."),
      ("", "Verdict: good at 30–50 cm — stacked-pile morphology.", True)],
     "Torres & Fraternali 2023"),

    ("9 · Plastic", "material (spectrum)",
     [("", "Satellite: morphological only at 30–50 cm (WV-3 spectral anomaly; Pléiades + S-2 index). True spectral plastic only at Sentinel-2 10 m (out of scope)."),
      ("", "Drone: YES — UAV-SWIR Attention-U-Net 96.8% acc / 91.1% F1 (2026); UAV-RGB + IoT 92% on rivers (2025)."),
      ("", "Method: index/anomaly at VHR; deep segmentation on UAV-SWIR for polymer type."),
      ("", "Verdict: presence/extent at VHR; polymer identity needs SWIR → drone-HSI.", True)],
     "UAV-SWIR U-Net RS 18:182, 2026 · UAV+IoT J.Haz.Mat.Adv. 2025"),

    ("10 · Big bags (FIBC)", "object (shape), small",
     [("", "Satellite: marginal — only AerialWaste (50 labels); a ~1 m³ bag ≈ 3×3 px at 30 cm."),
      ("", "Drone: NO dedicated detector."),
      ("", "Method: none specific."),
      ("", "Verdict: near-GAP — object but at the resolution limit + data-scarce.", True)],
     "Torres & Fraternali 2023"),

    ("11 · Tanks / cisterns", "object (shape)",
     [("", "Satellite: YES, mature — YOLOv7-OT 90% acc / 95.9% prec (2024); Ramachandran et al. Precision 0.962 / Recall 0.968, >169k tanks (Nature Comm. 2024)."),
      ("", "Drone: not specifically needed."),
      ("", "Method: object detection (circular tank tops)."),
      ("", "Verdict: excellent — and NOT a gap (correction vs first matrix).", True)],
     "YOLOv7-OT RS 16:4510, 2024 · Ramachandran et al. Nature Communications 2024"),

    ("12 · Tires", "object + spectral",
     [("", "Satellite: YES for piles — TIRe model on QuickBird 0.6 m (piles ≥100–400 tires); Disaitek qualifies tires at 30 cm; AerialWaste 45 labels."),
      ("", "Drone: partial."),
      ("", "Method: reflectance/decision-tree + object detection."),
      ("", "Verdict: good for piles; dark targets confuse with shadow/water.", True)],
     "TIRe model (CIWMB/NASA NTRS) · Disaitek/Airbus 2024 · Torres & Fraternali 2023"),

    ("13 · Asbestos-cement roofing", "material (spectrum)",
     [("", "Satellite + SWIR: Bonifazi 2026 WV-3 VNIR+SWIR, MLC building-level (Mantua, open Python); Saba WV-3 8-VNIR, Macro-F1 97.6%; EnMAP 30 m ACE 91.4% (out of scope = SWIR upper bound)."),
      ("", "Satellite/aerial, no SWIR: Abbasi 2024 Nearmap aerial, DenseNet+LSTM multi-temporal, OA 95.8–96.0%, AC 94%."),
      ("", "Drone: asbestos-slate drone-RGB DL training data (2023)."),
      ("", "Verdict: corrugated shape detectable at 30–50 cm; material confirmation needs SWIR (≥1.2 m WV-3) or very fine VHR + temporal.", True)],
     "Bonifazi Geomatics 6:41, 2026 · Saba 2026 · Abbasi RSASE 2024 · Shepherd Sci.Rep. 2025"),
]

for ttl, tag, items, ref in CLASSES:
    s = new()
    title(s, ttl, tag)
    bullets(s, items, top=Inches(1.75), size=16, gap=9)
    footer(s, ref)

# ---- Summary matrix ----
s = new()
title(s, "All 13 classes at a glance", "S = satellite · D = drone · ✓ served · ◐ partial / morphology-only · ✗ gap")
rows = [
    ["1 Rubble / C&D", "✓", "✓", "heap morphology — strong"],
    ["2 Foundry slag", "✗", "✗", "GAP — spectral + data-starved"],
    ["3 Vehicles / ELV", "✓", "✓", "object — excellent (ELV sub-gap)"],
    ["4 Scrap metal", "◐", "✗", "scenes yes; composition GAP"],
    ["5 Bulky items", "✓", "◐", "object — good"],
    ["6 Containers", "✓", "—", "object — strong"],
    ["7 Sludge", "✗", "✗", "GAP — context only"],
    ["8 Wood / firewood", "✓", "◐", "pile morphology — good"],
    ["9 Plastic", "◐", "✓", "presence at VHR; type needs SWIR"],
    ["10 Big bags", "◐", "✗", "near-GAP — small + scarce"],
    ["11 Tanks / cisterns", "✓", "—", "object — excellent"],
    ["12 Tires", "✓", "◐", "piles — good (dark-target)"],
    ["13 Asbestos", "✓", "✓", "shape yes; material needs SWIR"],
]
table(s, ["Class", "S", "D", "Verdict at 30–50 cm"],
      rows, [Inches(3.0), Inches(0.7), Inches(0.7), Inches(7.5)], Inches(1.7), fsize=11)

# ---- Gaps ----
s = new()
title(s, "The gaps — and the thesis space")
bullets(s, [
    "No dedicated VHR-satellite or drone detector for: foundry slag (2), sludge (7), big bags (10), scrap composition (4).",
    "Asbestos at true 30–50 cm is not validated — strong results use WV-3 ≥1.2 m with SWIR, EnMAP 30 m, or ≤25 cm aerial.",
    "Plastic polymer type needs SWIR — at 30–50 cm satellite you get presence/extent only.",
    "The whole field is binary / aggregate — no public per-class (13-material) detector at the VHR satellite operating point.",
    ("", "This empty cell is the thesis contribution space.", True),
    "Correction vs first matrix: tanks (11) are well-served, NOT a gap.",
], size=17)

prs.save("/home/alepot55/Desktop/uni/Tesi/vhr_13classes_deck.pptx")
print("saved vhr_13classes_deck.pptx —", len(prs.slides._sldIdLst), "slides")
