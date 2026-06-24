#!/usr/bin/env python3
"""Minimal B/W brief deck: asbestos-roof signature extraction (7 slides).

Style: white background, black text, thin rule under titles, page numbers.
Figures: slide-ready variants from slides/figs/ (no internal titles).
Output: asbestos/slides/asbestos_signatures_brief.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "slides" / "figs"
OUT = ROOT / "slides"
OUT.mkdir(exist_ok=True)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)
FONT = "Helvetica"
BLACK = RGBColor(0x11, 0x11, 0x11)
GREY = RGBColor(0x77, 0x77, 0x77)
RULE = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
N_PAGE = 0


def text(slide, x, y, w, h, runs, size, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, leading=1.0, space_after=0):
    """runs: str (multi-line) or list of (text, bold) paragraph tuples."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    items = runs.split("\n") if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        t, b = (item, bold) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = b
        p.font.color.rgb = color
        p.alignment = align
        p.line_spacing = leading
        p.space_after = Pt(space_after)
    return box


def rule(slide, y, x0=MARGIN, x1=SLIDE_W - MARGIN, weight=1.2, color=RULE):
    ln = slide.shapes.add_connector(1, x0, y, x1, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)


def header(slide, title_text):
    text(slide, MARGIN, Inches(0.42), SLIDE_W - 2 * MARGIN, Inches(0.7),
         title_text, 23, bold=True)
    rule(slide, Inches(1.18))


def footer(slide, note=""):
    global N_PAGE
    N_PAGE += 1
    if note:
        text(slide, MARGIN, SLIDE_H - Inches(0.52), SLIDE_W - 2 * MARGIN - Inches(0.6),
             Inches(0.4), note, 10.5, color=GREY)
    text(slide, SLIDE_W - MARGIN - Inches(0.5), SLIDE_H - Inches(0.52),
         Inches(0.5), Inches(0.4), str(N_PAGE), 10.5, color=GREY, align=PP_ALIGN.RIGHT)


def picture(slide, path: Path, top_in: float, bottom_in: float = 0.75):
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = Inches(SLIDE_H.inches - top_in - bottom_in)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((SLIDE_W - w) / 2), Inches(top_in),
                             width=w, height=h)


def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── 1 · Title ────────────────────────────────────────────────────────────────
s = slide()
text(s, MARGIN, Inches(2.05), SLIDE_W - 2 * MARGIN, Inches(0.4),
     "PERIVALLON · ASBESTOS PILOT", 13, color=GREY)
rule(s, Inches(2.55), x1=MARGIN + Inches(2.2))
text(s, MARGIN, Inches(2.75), SLIDE_W - 2 * MARGIN, Inches(1.6),
     "Asbestos roofs: first real spectral signatures\nfrom SuperDove imagery", 32,
     bold=True, leading=1.05)
text(s, MARGIN, Inches(4.55), SLIDE_W - 2 * MARGIN, Inches(1.0),
     "Regional asbestos ground truth (Mappatura 2020) × Planet PSScene 8-band, 12 May 2026\n"
     "Alessandro Potenza — 12 June 2026", 15, color=GREY, leading=1.25)
footer(s)

# ── 2 · Data / where ─────────────────────────────────────────────────────────
s = slide()
header(s, "The PSScene strips cover 1,096 mapped roofs — 966 usable signatures")
picture(s, FIGS / "map.png", 1.5)
footer(s, "Ground truth: Regione Lombardia “Mappatura 2020”, 10.9k roofs · usable = ≥ 150 m², "
          "≥ 5 clear pixels · Brescia province (Manerbio, Leno, Bagnolo Mella) · 5 strips, single day.")

# ── 3 · Method ───────────────────────────────────────────────────────────────
s = slide()
header(s, "Extraction pipeline: from GT polygons to one signature per roof")
steps = [
    ("GT selection", "Mappatura 2020 census (most recent regional mapping), area ≥ 150 m²"),
    ("Edge erosion", "polygon buffered −3 m inward: at 3 m GSD this removes mixed roof/ground "
                     "border pixels, keeping pure roof pixels only"),
    ("Pixel quality", "Planet udm2 clear mask; a roof needs ≥ 5 valid pixels"),
    ("Signature", "per-band median of surface reflectance — 8 SuperDove bands (Coastal Blue → NIR); "
                  "overlapping strips → keep the observation with more clear pixels"),
    ("Vegetation screening", "NDVI < 0.2 mineral (bare roof) · NDVI ≥ 0.3 vegetated "
                             "(likely remediated or overgrown since the 2020 mapping)"),
    ("Similarity", "SAM — angle between signatures as vectors: compares spectral shape, "
                   "ignores overall brightness"),
]
y = 1.55
for i, (lead, body) in enumerate(steps, 1):
    text(s, MARGIN, Inches(y), Inches(0.45), Inches(0.4), f"{i}", 15, bold=True, color=GREY)
    text(s, MARGIN + Inches(0.5), Inches(y), Inches(2.5), Inches(0.4), lead, 15, bold=True)
    text(s, MARGIN + Inches(3.2), Inches(y), SLIDE_W - 2 * MARGIN - Inches(3.2),
         Inches(0.85), body, 14, leading=1.05)
    y += 0.88
footer(s, "966 of 1,096 covered roofs pass all filters · surface reflectance, scale 10⁻⁴.")

# ── 4 · Visual check ─────────────────────────────────────────────────────────
s = slide()
header(s, "Extraction check on the imagery: GT outlines match the buildings")
picture(s, FIGS / "chips.png", 1.5)
footer(s, "Planet RGB, 3 m, GT outline in yellow · row 1: largest mineral roofs — industrial sheds, "
          "classic asbestos-cement · row 2: roofs flagged as vegetated by NDVI.")

# ── 5 · Result ───────────────────────────────────────────────────────────────
s = slide()
header(s, "Bare roofs match the lab reference — NDVI flags likely-remediated roofs")
picture(s, FIGS / "fan.png", 1.5)
footer(s, "GT is from 2020, image from 2026: the vegetated subset (red-edge shape) marks roofs likely "
          "remediated or overgrown in between · mineral subset SAM 6.6° vs USGS asbestos-cement model "
          "resampled to the 8 SuperDove bands.")

# ── 6 · Why SWIR matters ─────────────────────────────────────────────────────
s = slide()
header(s, "Why SWIR matters: only asbestos-cement shows the 2.31 µm Mg-OH feature")
picture(s, FIGS / "band_depth.png", 1.5)
footer(s, "Continuum-removed band depth, USGS splib07a; asbestos-cement modelled as chrysotile–concrete "
          "mix (Cilia 2015) · the feature grows with weathering — measurable by WV-3 S8, washed out by "
          "Sentinel-2 B12 (FWHM 175 nm), invisible to SuperDove.")

# ── 7 · Takeaway ─────────────────────────────────────────────────────────────
s = slide()
header(s, "What the VNIR bands can and cannot do")
points = [
    ("Asbestos-cement vs terracotta: separable",
     "SAM 21.6° on SuperDove bands — the common Lombardy roof is not a confuser"),
    ("Asbestos-cement vs bare concrete: not separable",
     "SAM ~4° on every band set — the hard confuser stays hard in VNIR"),
    ("The diagnostic feature lives in the SWIR",
     "Mg-OH absorption at 2.31 µm: needs WV-3; out of reach for SuperDove and Sentinel-2"),
    ("Next step",
     "extend the extraction to more covered roofs and validate the SWIR feature "
     "on the project WV-3 scenes"),
]
y = 1.75
for lead, body in points:
    text(s, MARGIN, Inches(y), SLIDE_W - 2 * MARGIN, Inches(0.4), lead, 17, bold=True)
    text(s, MARGIN, Inches(y + 0.42), SLIDE_W - 2 * MARGIN, Inches(0.5), body, 14, color=GREY)
    y += 1.22
footer(s)

out = OUT / "asbestos_signatures_brief.pptx"
prs.save(out)
print("Deck:", out)
