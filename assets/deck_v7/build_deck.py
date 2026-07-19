"""Deck v7 complete. 25 slides, full in-scope SOTA coverage, plain block
diagrams, real paper figures. Content per Thomas's review."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT=os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIG=os.path.join(ROOT,"assets/deck_v7/figs")
OUTP=os.path.join(ROOT,"assets/deck_v7/deck_v7.pptx")
INK=RGBColor(0x1A,0x1A,0x1A); GREY=RGBColor(0x66,0x66,0x66)
HDR=RGBColor(0xDD,0xDD,0xDD); ALT=RGBColor(0xF5,0xF5,0xF5); WHITE=RGBColor(0xFF,0xFF,0xFF)
F="Calibri"
prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
BLANK=prs.slide_layouts[6]
def IN(v): return Inches(v)
def tf_fill(tf,items,wrap=True):
    tf.clear(); tf.word_wrap=wrap
    for k,it in enumerate(items):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(it.get("sa",6)); p.line_spacing=it.get("ls",1.12)
        r=p.add_run(); r.text=it["t"]
        r.font.name=F; r.font.size=Pt(it.get("sz",14)); r.font.bold=it.get("b",False)
        r.font.italic=it.get("i",False); r.font.color.rgb=it.get("c",INK)
def slide(): return prs.slides.add_slide(BLANK)
def title(s,t):
    tb=s.shapes.add_textbox(IN(0.45),IN(0.28),IN(9.10),IN(0.6))
    tf_fill(tb.text_frame,[{"t":t,"b":True,"sz":22}])
def body(s,items,top=1.15,left=0.45,w=9.10,h=3.85,mid=False):
    tb=s.shapes.add_textbox(IN(left),IN(top),IN(w),IN(h)); tf_fill(tb.text_frame,items)
    if mid: tb.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return tb
def foot(s,t):
    tb=s.shapes.add_textbox(IN(0.45),IN(5.26),IN(9.10),IN(0.26))
    tf_fill(tb.text_frame,[{"t":t,"sz":8,"c":GREY,"sa":0}])
def img(s,path,left,top,w,h):
    from PIL import Image
    iw,ih=Image.open(path).size; ar=iw/ih
    W=IN(w); H=int(W/ar)
    if H>IN(h): H=IN(h); W=int(H*ar)
    s.shapes.add_picture(path,IN(left)+(IN(w)-W)//2,IN(top)+(IN(h)-H)//2,width=W,height=H)
def table(s,data,left=0.45,top=1.25,w=9.10,colw=None,fs=9.5,rh=0.46,hh=0.38,vc=False):
    nr=len(data); nc=len(data[0]); th=hh+(nr-1)*rh
    if vc: top=1.12+max(0,(4.0-th)/2.0)
    gt=s.shapes.add_table(nr,nc,IN(left),IN(top),IN(w),IN(th)).table
    if colw:
        for j,cw in enumerate(colw): gt.columns[j].width=IN(cw)
    gt.rows[0].height=IN(hh)
    for i in range(1,nr): gt.rows[i].height=IN(rh)
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=IN(0.05); c.margin_right=IN(0.04); c.margin_top=IN(0.015); c.margin_bottom=IN(0.015)
            c.fill.solid(); c.fill.fore_color.rgb=HDR if i==0 else (ALT if i%2==0 else WHITE)
            for p in c.text_frame.paragraphs:
                p.line_spacing=1.0
                for r in p.runs:
                    r.font.name=F; r.font.size=Pt(fs+0.5 if i==0 else fs)
                    r.font.bold=(i==0); r.font.color.rgb=INK
    return gt

# 1 ── title
s=slide()
tb=s.shapes.add_textbox(IN(0.7),IN(1.7),IN(8.6),IN(1.4))
tf_fill(tb.text_frame,[{"t":"Classification of waste materials in very-high-resolution satellite imagery","b":True,"sz":30}])
tb2=s.shapes.add_textbox(IN(0.7),IN(3.15),IN(8.6),IN(0.5))
tf_fill(tb2.text_frame,[{"t":"State of the art and thesis proposal","sz":17,"i":True,"c":GREY}])
tb3=s.shapes.add_textbox(IN(0.7),IN(4.5),IN(8.6),IN(0.4))
tf_fill(tb3.text_frame,[{"t":"Alessandro Potenza  ·  M.Sc. Computer Science and Engineering, AI  ·  PERIVALLON","sz":12,"c":GREY}])

# 2 ── context
s=slide(); title(s,"Context")
body(s,[
 {"t":"Illegal waste dumping is an environmental crime with direct public-health consequences. Agencies (ARPA) have limited inspection capacity, and priority depends on what is dumped: rubble, plastics and asbestos-cement imply very different hazards.","sz":12.5},
 {"t":"Detecting dump sites from images is mature. Recognising the material is not: the reference survey (50 works, 1987-2023, almost all RGB) marks it as an open problem.","sz":12.5},
 {"t":"One thesis in the group has addressed material classification directly (Alari 2024). It is examined here alongside the rest of the literature.","sz":12.5},
],top=1.00,h=1.85,mid=False)
img(s,os.path.join(FIG,"torres_examples.png"),0.7,2.95,8.6,2.20)
foot(s,"Images: AerialWaste positive examples, Torres and Fraternali 2023 (CC BY 4.0)  ·  survey: Fraternali et al. 2024")

# 3 ── task definition
s=slide(); title(s,"Task definition")
body(s,[
 {"t":"Objective: given a VHR satellite image of a suspected dump site, output the set of waste materials present.","b":True,"sz":13,"sa":14},
 {"t":"Task: multi-label classification at image level.","sz":13,"sa":14},
 {"t":"Input: VHR optical imagery, GSD 0.2 to 1.3 m. Aerial RGB for the baseline (AerialWaste); satellite VNIR for the multispectral arm, up to 8 bands. SWIR excluded: not in the planned acquisitions, and its 3.7 m GSD does not match the task. Sentinel-2 excluded: too coarse (10-20 m).","sz":13,"sa":14},
 {"t":"Technique: classification, not object detection or segmentation. Annotations are image-level (Alari 2024: 11,477 multi-label); waste piles have no stable shape for boxes; no segmentation masks exist for these datasets.","sz":13,"sa":14},
 {"t":"Research question: does VNIR input improve waste-material classification over RGB, and for which materials?","b":True,"sz":13},
],top=1.15,h=3.95)
foot(s,"Annotations: Alari 2024  ·  imagery: WorldView-3, Pléiades Neo (VNIR + pan)")

# 4 ── task scheme
s=slide(); title(s,"Task scheme")
img(s,os.path.join(FIG,"task_io.png"),0.7,1.35,8.6,3.1)
body(s,[{"t":"One image in, a set of material labels out. The same scheme runs with RGB or VNIR input; only the first layer of the network changes.","sz":12}],top=4.55,h=0.6,mid=False)
foot(s,"Tile: Alari 2024. The multi-label formulation follows the group predecessor.")

# 5 ── materials taxonomy (full 10-class samples)
s=slide(); title(s,"The material taxonomy")
img(s,os.path.join(FIG,"alari_row1.png"),0.85,1.10,8.3,1.75)
img(s,os.path.join(FIG,"alari_row2.png"),0.85,2.95,8.3,1.75)
body(s,[{"t":"Ten-category grouping used in the group predecessor; the full annotation set counts 13 categories (adds foundry waste, sludge, tanks as separate classes).","sz":11.5,"c":GREY}],top=4.72,h=0.5,mid=False)
foot(s,"Samples and taxonomy: Alari 2024 (politesi 10589/230633), built on AerialWaste annotations")

# 6 ── materials decision
s=slide(); title(s,"Materials: which are considered and why")
table(s,[
 ["Material","Target","Spectral analysis","Reason"],
 ["Rubble, plastic, wood, tires","yes","yes","frequent and confusable in RGB"],
 ["Asbestos-cement","yes","yes, dedicated pilot","public regional ground truth (WFS)"],
 ["Vehicles, tanks, containers, scrap, bulky, big bags","yes","no","shape-based; RGB sufficient in literature"],
 ["Sludge, foundry waste","yes","no","few labels; visually ambiguous at this GSD"],
],left=0.45,top=1.30,w=9.10,colw=[3.4,0.95,1.85,2.9],fs=9.5,rh=0.55,hh=0.42)
body(s,[{"t":"All 13 categories remain classification targets, for continuity with the group dataset. The RGB-versus-VNIR analysis concentrates on the subset where colour is ambiguous and labels are sufficient.","sz":12}],top=4.20,h=0.85)
foot(s,"Decision criteria: visual ambiguity in RGB, label availability, hazard relevance")

# 7 ── material x literature coverage
s=slide(); title(s,"Per-material coverage in the literature")
table(s,[
 ["Material","Dedicated works at task scale","Notes"],
 ["Asbestos-cement","Saba 2026, Bonifazi 2026, Abbasi 2024, Cilia 2015","roofs only, never inside a waste taxonomy"],
 ["Tanks","Ramachandran 2024, YOLOv7-OT 2024","mature object detection"],
 ["Vehicles","vehicle-detection DA 2020","VHR object detection"],
 ["Containers","truck-and-container 2025","size and status classification"],
 ["Scrap","ELV Hybrid-YOLOv5 2025","close-range infrared, not satellite"],
 ["Rubble / inert","CWLD 2024; debris volume UAV 2022","C&D segmentation; UAV photogrammetry"],
 ["Plastic, wood, tires, bulky, big bags","none found in scope","only as classes in AerialWaste / Alari"],
 ["Sludge, foundry waste","none found","only as classes in Alari"],
],left=0.45,top=1.22,w=9.10,colw=[2.30,3.55,3.25],fs=8.5,rh=0.43,hh=0.36)
foot(s,"Scope: terrestrial, GSD compatible with the task. This table answers how thin the material-level literature actually is.")

# 8 ── search method
s=slide(); title(s,"Literature search: method and numbers")
body(s,[
 {"t":"Scripted queries on the Scopus API (waste detection in remote sensing; asbestos roof mapping), deduplication, manual screening with explicit criteria, snowballing from the Fraternali 2024 survey and the Alari 2024 references.","sz":12.5},
],top=1.05,h=1.0,mid=False)
img(s,os.path.join(FIG,"search_flow.png"),0.55,2.15,8.9,2.55)
foot(s,"Search artifacts: papers/literature_search (queries, raw and deduplicated records)  ·  library: papers/notes")

# 9 ── kept vs excluded
s=slide(); title(s,"What was kept, what was excluded")
table(s,[
 ["Excluded group","Examples","Reason"],
 ["Sentinel-2 and marine debris","MARIDA 2022, Tisza 2023","10-20 m: pixels too mixed; different domain"],
 ["Spaceborne hyperspectral","Shepherd 2025 (EnMAP), EMIT 2025","material evidence, but 30-60 m GSD"],
 ["SWIR-based VHR","Aguilar 2021, Aguilar 2025, Zhou 2021","SWIR not in our acquisitions; 3.7 m GSD"],
 ["Laboratory / conveyor","Vitek 2025, SpectralWaste 2024","not EO; Vitek reused as band-selection evidence"],
 ["EO foundation models","DOFA 2024, AnySat 2025, Prithvi 2024, +6","pretrained at 10-30 m; sub-metre transfer unproven"],
],left=0.45,top=1.30,w=9.10,colw=[2.55,3.20,3.35],fs=9,rh=0.50,hh=0.38)
body(s,[{"t":"Kept: works on terrestrial waste or roof materials at task-compatible GSD, plus the reference spectral library. The excluded papers stay in the annotated library as context.","sz":11.5}],top=4.35,h=0.7,mid=False)
foot(s,"Full accounting in papers/INDEX.md (47 papers, each with a structured note)")

# 10 ── site-level table
s=slide(); title(s,"Related work: site-level waste detection")
table(s,[
 ["Work (year)","Input / GSD","Task","Method","Result"],
 ["Gibellini 2025","aerial RGB, 20 cm","classification","Swin-T, RSP pretraining","F1 92.0; cross-region -5.1"],
 ["AW ensemble 2025 *","aerial RGB","classification","CNN + transformer ensemble","binary F1 92.4"],
 ["Disaitek 2024","Pléiades Neo, 30 cm","detection","operational service","~95% on sites >2 m2 (vendor)"],
 ["CascadeDumpNet 2024","Pléiades, 50 cm","object detection","two-stage CNN + AutoML","mAP 84.6, cross-city transfer"],
 ["Sun 2023","VHR RGB, 0.3-1 m","object detection","BCA-Net (Faster R-CNN)","~2,500 dumpsites, 28 cities"],
 ["CWLD 2024","GF-2 + GE, 0.5-0.8 m","segmentation","DeepLabV3+ variant","F1 88.9, construction waste"],
 ["AerialWaste, Torres 2023","aerial RGB, 20-50 cm","dataset + cls.","ResNet-FPN baseline","F1 80.7; 22 material tags"],
],colw=[2.05,1.85,1.45,1.95,1.80],fs=8.5,rh=0.42,hh=0.36,vc=True)
foot(s,"All RGB: they answer where a site is, not what it contains. * preprint; Disaitek is a vendor figure.")

# 11 ── deep-dive AerialWaste
s=slide(); title(s,"AerialWaste: the reference dataset")
body(s,[
 {"t":"10,434 locations from ARPA Lombardia records, 487 municipalities; three sources: AGEA orthophotos (20 cm), WorldView-3 (30 cm), Google Earth (50 cm).","sz":11.5,"sa":10},
 {"t":"Binary site labels for detection, plus 22 material tags (type of visible object + storage mode) annotated by experts.","sz":11.5,"sa":10},
 {"t":"Material tags cover about 72% of positives, but the dataset was published for detection: material tags are metadata, not a benchmark.","sz":11.5,"sa":10},
 {"t":"This is the label base both Gibellini 2025 and Alari 2024 build on.","sz":11.5,"sa":10},
],top=1.15,left=0.45,w=4.10,h=3.90)
img(s,os.path.join(FIG,"aerialwaste_grid.png"),4.75,1.10,4.85,4.00)
foot(s,"Image: annotated-label examples, Torres and Fraternali 2023, Scientific Data (CC BY 4.0)")

# 12 ── deep-dive Gibellini
s=slide(); title(s,"Gibellini 2025: the site-level baseline")
body(s,[
 {"t":"Binary waste / no-waste classification on AerialWaste; Swin-T with remote-sensing pretraining, two-step fine-tuning.","sz":11.5,"sa":10},
 {"t":"F1 92.02 in domain. Cross-region generalisation drops 5.1 points on average (Greece 85.4, Serbia 83.8, Romania 91.5).","sz":11.5,"sa":10},
 {"t":"Saliency maps confirm the model looks at waste heaps, but the output is presence only: no material information.","sz":11.5,"sa":10},
 {"t":"This architecture and training recipe is the starting point of the proposal.","sz":11.5,"sa":10},
],top=1.15,left=0.45,w=4.10,h=3.90)
img(s,os.path.join(FIG,"gibellini_cams.png"),4.75,1.30,4.85,3.55)
foot(s,"Image: true positives with saliency overlays, Gibellini et al. 2025 (Greek generalisation set)")

# 13 ── deep-dive Alari
s=slide(); title(s,"Alari 2024: the closest work")
body(s,[
 {"t":"First work in the group to frame waste-material recognition as multi-label classification. It entered the library through the same screening as every other paper; the coverage table shows it is the only in-scope multi-material work, which makes it the operational comparison.","sz":12,"sa":14},
 {"t":"Dataset: 11,477 multi-label annotations over 13 categories; 3,190 positive and 7,190 negative images, built on AerialWaste imagery and ARPA records.","sz":12,"sa":14},
 {"t":"Models: ResNet-50 and Swin backbones with FPN; three classification-head designs; weighted binary cross-entropy for label imbalance; different pretraining sources compared.","sz":12,"sa":14},
 {"t":"Results: weighted F1 69.21 on five categories, 59.42 on ten. Growing the taxonomy from 5 to 10 costs 9.8 points.","sz":12,"sa":14},
 {"t":"Input is RGB only. The spectral dimension is untouched: that is the opening this thesis takes.","b":True,"sz":12},
],top=1.20,h=3.90)
foot(s,"Alari 2024, M.Sc. thesis, PoliMi (advisor Fraternali, co-advisor Gibellini), politesi 10589/230633")

# 13b -- Alari per-material results
s=slide(); title(s,"Alari 2024: results per material")
img(s,os.path.join(FIG,"alari_perclass.png"),0.40,1.15,5.55,3.90)
body(s,[
 {"t":"Per-class F1 of the ten-category model. The weighted average, 59.4, hides a split: classes defined by shape or extent stay near 70; classes defined by material appearance fall below 45.","sz":11.5,"sa":12},
 {"t":"The thesis attributes the low scores to few annotations (tires, big bags), high intra-class variance (wood) and inter-class similarity in RGB.","sz":11.5,"sa":12},
 {"t":"Three of the four spectral-analysis targets sit in the bottom half. Rubble is the exception: large extents make it visible even in RGB.","sz":11.5,"sa":12},
 {"t":"This per-class picture is what the band ablation is designed to move.","b":True,"sz":11.5},
],top=1.20,left=6.15,w=3.45,h=3.85)
foot(s,"Data: Alari 2024, Table 4.13 (ResNet-50 + IDA, ten categories), politesi 10589/230633")

# 14b -- what Alari's conclusions leave open
s=slide(); title(s,"Alari 2024: what the conclusions leave open")
body(s,[
 {"t":"1.  Dataset extension: 5 of the 13 classes have fewer than 400 samples; underrepresented classes limit generalisation.","sz":12.5,"sa":13},
 {"t":"2.  Multi- and hyper-spectral imagery: a wider spectral input \u201cmay also prove useful in the identification of waste types\u201d. Named explicitly as a future direction.","sz":12.5,"sa":13},
 {"t":"3.  Synthetic augmentation: high co-occurrence between classes confounds the model; single-category images would isolate class representations.","sz":12.5,"sa":13},
 {"t":"4.  Multi-class pre-training: fewer than 300 single-annotation images make it infeasible on the current dataset.","sz":12.5,"sa":13},
 {"t":"5.  Object detection: boxes could focus attention, at the cost of new annotations.","sz":12.5,"sa":13},
 {"t":"Point 2 is this thesis. Points 1 and 3 explain part of the per-class collapse on the previous slide.","b":True,"sz":12.5},
],top=1.10,h=4.0)
foot(s,"Alari 2024, Chapter 5, Conclusions and Future Works. The thesis is compared like every other library paper.")



# 14 ── material-level table
s=slide(); title(s,"Related work: material-level classification")
table(s,[
 ["Work (year)","Input / GSD","Task","Result"],
 ["Alari 2024 (PoliMi)","satellite RGB","multi-label cls.","wF1 69.2 (5 cat.), 59.4 (10 cat.)"],
 ["Saba 2026 *","WV-3 VNIR, 1.24 m","pixel cls.","asbestos, Macro-F1 97.6"],
 ["Bonifazi 2026","WV-3 VNIR+SWIR","pixel cls.","asbestos roofs, removal tracking"],
 ["Abbasi 2024 *","aerial RGB","OBIA cls.","asbestos, OA ~96 (shape + time)"],
 ["Cilia 2015","airborne HSI, 3 m","pixel cls.","asbestos, PA 89 / UA 86"],
],left=0.45,top=1.35,w=9.10,colw=[2.35,2.20,1.85,2.70],fs=9.5,rh=0.50,hh=0.40)
body(s,[{"t":"One multi-material predecessor; everything else addresses a single material, asbestos, on roofs. No work measures RGB versus VNIR on waste materials.","sz":12}],top=4.30,h=0.75,mid=False)
foot(s,"* paywalled or preprint, reported as such")

# 15 ── deep-dive asbestos
s=slide(); title(s,"The asbestos line, in detail")
body(s,[
 {"t":"Saba 2026: WorldView-3, VNIR only, 32 classifiers compared; best Macro-F1 97.6 per-pixel. Red Edge and NIR carry the discrimination.","sz":11.5,"sa":10},
 {"t":"Bonifazi 2026: WorldView-3 multi-temporal workflow; building-level decisions from pixel classification; tracks roof removals between acquisitions.","sz":11.5,"sa":10},
 {"t":"Abbasi 2024: aerial RGB with temporal features, OA ~96: shape and time can substitute spectra at fine GSD.","sz":11.5,"sa":10},
 {"t":"Cilia 2015: airborne hyperspectral, PA 89 / UA 86, plus a weathering index from red/NIR: the degradation angle.","sz":11.5,"sa":10},
 {"t":"Asbestos slate on drone RGB (2023): partially recognisable by shape and texture alone.","sz":11.5,"sa":10},
 {"t":"Why it matters here: public ground truth (Lombardy WFS, 10,903 roofs) and VNIR evidence make asbestos the natural pilot.","b":True,"sz":11.5,"sa":10},
],top=1.10,left=0.45,w=4.10,h=4.0)
img(s,os.path.join(FIG,"bonifazi_example.png"),4.85,1.20,4.65,3.85)
foot(s,"Image: building-level asbestos classification on WorldView-3, Bonifazi et al. 2026, Geomatics (CC BY 4.0)")

# 15b -- asbestos weathering in VNIR (Cilia 2015)
s=slide(); title(s,"Asbestos weathering is visible in VNIR")
body(s,[
 {"t":"Cilia 2015 measured asbestos-cement roofs of different age and exposure, on the ground and from the MIVIS airborne sensor. The 0.48-0.82 um window is shown: pure VNIR.","sz":11.5,"sa":12},
 {"t":"Older roofs are darker across the whole visible range; vegetation colonising the surface adds a chlorophyll absorption near 680 nm; north-exposed roofs amplify both effects.","sz":11.5,"sa":12},
 {"t":"From these features the paper builds a deterioration index and ranks removal priority; the effects of age and exposure are statistically significant (ANOVA, p < 0.001).","sz":11.5,"sa":12},
 {"t":"The discriminating features fall inside the VNIR window of the planned sensors: degradation state, not only presence, is within reach of the pilot.","b":True,"sz":11.5},
],top=1.15,left=0.45,w=4.55,h=3.95)
img(s,os.path.join(FIG,"cilia_weathering.png"),5.25,1.10,4.30,4.00)
foot(s,"Figure: field and MIVIS spectra of asbestos-cement roofs, Cilia et al. 2015, ISPRS IJGI 4(2) (CC BY)")


# 17 ── rgb limits + vnir
s=slide(); title(s,"Where RGB falls short, and what VNIR adds")
body(s,[
 {"t":"Different materials share the same colour at 0.3-1.3 m: plastic sheets, asbestos-cement and concrete all appear grey.","sz":12,"sa":11},
 {"t":"In the group predecessor, moving from 5 to 10 material categories costs 9.8 F1 points (69.2 to 59.4). Finer material distinctions are the hard part.","sz":12,"sa":11},
 {"t":"Red Edge and NIR separate vegetation, bare soil and weathered surfaces; in Saba 2026 they drive asbestos discrimination.","sz":12,"sa":11},
 {"t":"Whether this helps waste materials, and which ones, has not been measured. It is the object of this thesis.","sz":12},
],top=1.15,left=0.45,w=3.85,h=3.90)
img(s,os.path.join(FIG,"vnir_signatures.png"),4.45,1.20,5.15,3.85)
foot(s,"Reflectance: USGS splib07a (Kokaly et al. 2017), 400-1050 nm. Band centres: Maxar, Airbus.")

# 17b -- how many extra bands close the gap
s=slide(); title(s,"How many extra bands close the RGB gap")
body(s,[
 {"t":"Controlled-condition evidence (Vitek 2025, laboratory hyperspectral, 10 construction and demolition materials): RGB alone reaches 0.87 accuracy; two well-chosen extra bands bring it to 0.96; the full 768-band spectrum adds almost nothing more.","sz":12,"sa":8},
 {"t":"The optimal extra bands fall at 650-750 and 850-1000 nm: where the Red Edge and NIR bands of WorldView-3 and Pléiades Neo sit.","sz":12},
],top=1.02,h=1.30)
img(s,os.path.join(FIG,"bands_evidence.png"),0.50,2.28,9.00,2.52)
body(s,[{"t":"Laboratory study, reported as such: it bounds what spectra can add. Whether it transfers to satellite GSD is what this thesis measures.","sz":10.5,"c":GREY}],top=4.82,h=0.40)
foot(s,"Vitek, Zbiral, Nezerka 2025. Resources, Conservation and Recycling. Band ranges: Maxar, Airbus data sheets")


# 18 ── imagery
s=slide(); title(s,"Available imagery")
img(s,os.path.join(FIG,"bands_chart.png"),0.55,1.05,8.9,2.85)
body(s,[
 {"t":"Panchromatic: 0.31 m (WorldView-3), 0.30 m (Pléiades Neo). Access: commercial, or free quota via ESA proposal.","sz":12},
 {"t":"This is the band budget of the study: no SWIR, no Sentinel-2. RGB baseline runs on AerialWaste (20-50 cm).","sz":12},
],top=4.05,h=1.05,mid=False)
foot(s,"Band ranges: Maxar WorldView-3, Airbus Pléiades Neo data sheets")

# 19 ── synthesis
s=slide(); title(s,"The in-scope state of the art at a glance")
table(s,[
 ["Work","Input","Task","Material info"],
 ["Fraternali 2024 (survey)","-","survey, 50 works","declares the material gap"],
 ["AerialWaste 2023","aerial RGB","dataset","22 tags, metadata only"],
 ["Gibellini 2025","aerial RGB","classification","none (presence)"],
 ["Alari 2024","satellite RGB","multi-label cls.","13 categories, wF1 59-69"],
 ["AW ensemble 2025 *","aerial RGB","classification","none (binary)"],
 ["Sun 2023 / CascadeDumpNet 2024","VHR RGB","object detection","none (sites)"],
 ["Disaitek 2024 (vendor)","Pléiades Neo","detection","coarse type qualification"],
 ["CWLD 2024","GF-2 + GE","segmentation","one material (C&D)"],
 ["Saba / Bonifazi / Abbasi / Cilia","WV-3, aerial, HSI","pixel cls.","one material (asbestos)"],
 ["Tanks / vehicles / scrap works","VHR, close range","object detection","shape, not composition"],
],left=0.45,top=1.08,w=9.10,colw=[2.90,1.75,1.95,2.50],fs=8.5,rh=0.37,hh=0.33)
foot(s,"Ten lines summarise every in-scope line of work. Material-level evidence: one multi-label predecessor + one material studied in isolation.")

# 20 ── gaps
s=slide(); title(s,"What is missing in the literature")
body(s,[
 {"t":"1.  Multi-material classification has one direct precedent, with ample margin: wF1 59-69 (Alari 2024).","sa":20},
 {"t":"2.  No work measures the added value of VNIR bands over RGB for waste materials at very high resolution; Alari 2024 names the multispectral extension as future work.","sa":20},
 {"t":"3.  Results are reported as aggregate scores; per-material behaviour is not analysed.","sa":20},
 {"t":"4.  Generalisation across regions is rarely evaluated. At site level it costs 5 F1 points (Gibellini 2025).","sa":20},
 {"t":"5.  Asbestos is studied on roofs, in isolation, and never inside a waste-material taxonomy."},
],top=1.30,h=3.70)
foot(s,"Each gap maps to one element of the proposal on the next slides.")

# 21 ── proposal approach
s=slide(); title(s,"Proposed work: approach")
body(s,[
 {"t":"Multi-label image classification, continuing the group line (Gibellini 2025, Alari 2024). Same taxonomy, input extended from RGB to VNIR.","b":True,"sz":12},
 {"t":"Band ablation with everything else fixed: same architecture, same splits, only the input changes.","sz":12},
],top=1.05,h=1.10)
img(s,os.path.join(FIG,"pipeline.png"),0.55,2.15,8.9,2.75)
foot(s,"Backbone: Swin-T with remote-sensing pretraining (group baseline); extra bands enter via the input layer.")

# 22 ── pilot
s=slide(); title(s,"First step: the asbestos pilot")
body(s,[
 {"t":"1.  Extract roof polygons from the Lombardy WFS registry (10,903 mapped asbestos-cement roofs) and pair them with the available imagery.","sz":12.5,"sa":13},
 {"t":"2.  Build matched RGB and VNIR inputs for the same roofs; add negative roofs from regional building footprints.","sz":12.5,"sa":13},
 {"t":"3.  Train the same classifier on both inputs; compare F1, precision, recall on held-out areas.","sz":12.5,"sa":13},
 {"t":"4.  Decision point: the measured VNIR delta on a single, well-labelled material tells whether the full multi-label extension is worth the acquisition cost.","sz":12.5,"sa":13},
 {"t":"Why asbestos first: public pixel-accurate labels, one material, VNIR evidence in literature (Saba 2026), and direct regulatory value.","b":True,"sz":12.5},
],top=1.20,h=3.90)
foot(s,"Ground truth: Regione Lombardia WFS, Mappatura 2020 (EPSG:32632)")

# 23 ── evaluation
s=slide(); title(s,"Proposed work: evaluation")
body(s,[
 {"t":"Per-material F1 alongside weighted and macro averages: aggregates hide exactly the classes this thesis targets.","sz":12,"sa":11},
 {"t":"Delta versus the RGB baseline for each band configuration, with confidence intervals over repeated runs.","sz":12,"sa":11},
 {"t":"Generalisation: train and test on disjoint geographic areas, besides the standard split.","sz":12,"sa":11},
 {"t":"Reference points: wF1 69.2 / 59.4 (Alari 2024); Macro-F1 97.6 (Saba 2026, per-pixel) as upper reference for the pilot, not directly comparable.","sz":12,"sa":11},
 {"t":"If VNIR does not help a material, that is a documented negative result with practical value for sensor choice.","sz":12},
],top=1.15,left=0.45,w=4.7,h=3.9)
img(s,os.path.join(FIG,"eval_grid.png"),5.35,1.55,4.25,2.9)
foot(s,"Metrics: per-class F1, macro-F1, weighted F1, confusion matrices. Splits frozen before testing.")

# 24 ── references 1
s=slide(); title(s,"References (1/2)")
body(s,[
 {"t":"Alari 2024. Fighting environmental crime with deep learning: classifying waste materials from illegal landfills in satellite imagery. M.Sc. thesis, PoliMi, 10589/230633.","sz":10.5,"sa":12},
 {"t":"Fraternali et al. 2024. Solid waste detection, monitoring and mapping in remote sensing images: a survey. Waste Management / arXiv:2402.09066.","sz":10.5,"sa":12},
 {"t":"Gibellini et al. 2025. A deep learning pipeline for solid waste detection in remote sensing images. Waste Management Bulletin.","sz":10.5,"sa":12},
 {"t":"Torres, Fraternali 2023. AerialWaste: a dataset for illegal landfill discovery in aerial images. Scientific Data 10:63.","sz":10.5,"sa":12},
 {"t":"Sharmily et al. 2025. AerialWaste ensembles (preprint).  ·  Zhang, Ma 2024. CascadeDumpNet. Remote Sensing of Environment 313.","sz":10.5,"sa":12},
 {"t":"Sun et al. 2023. Revealing influencing factors on global waste distribution. Nature Communications 14:1444.","sz":10.5,"sa":12},
 {"t":"CWLD 2024. Construction waste landfill dataset. Scientific Data.  ·  Disaitek 2024, vendor report (Airbus).","sz":10.5,"sa":12},
 {"t":"Saba et al. 2026. J. Hazardous Materials.  ·  Bonifazi et al. 2026. Geomatics 6(3):41.  ·  Abbasi et al. 2024. RSASE.  ·  Cilia et al. 2015. ISPRS IJGI 4(2).","sz":10.5,"sa":12},
],top=1.20,h=3.90)
foot(s,"Full annotated library: 47 papers, papers/INDEX.md")

# 25 ── references 2
s=slide(); title(s,"References (2/2)")
body(s,[
 {"t":"Asbestos slate from drone imagery, 2023 (training-data study).","sz":10.5,"sa":12},
 {"t":"Ramachandran et al. 2024. Deep learning to map well pads and storage tanks. Nature Communications.","sz":10.5,"sa":12},
 {"t":"YOLOv7-OT 2024. Storage tank detection for large-scale remote sensing images. Remote Sensing.","sz":10.5,"sa":12},
 {"t":"Truck-and-container detection from satellite imagery, 2025 (preprint).","sz":10.5,"sa":12},
 {"t":"Vehicle detection with domain adaptation in VHR imagery, 2020.","sz":10.5,"sa":12},
 {"t":"ELV Hybrid-YOLOv5 2025. Non-ferrous metal detection in end-of-life vehicles (close-range infrared).","sz":10.5,"sa":12},
 {"t":"UAV solid-waste segmentation, 2024.  ·  C&D debris volume from UAV photogrammetry, 2022. Drones.","sz":10.5,"sa":12},
 {"t":"fCLIPSeg 2025. Event-agnostic debris segmentation.  ·  Kokaly et al. 2017. USGS Spectral Library v7 (splib07a). USGS DS 1035.","sz":10.5,"sa":12},
 {"t":"Vitek, Zbiral, Nezerka 2025. Critical wavelengths for construction and demolition waste materials. Resources, Conservation and Recycling.","sz":10.5,"sa":12},
],top=1.20,h=3.90)
foot(s,"24 works cited; the remaining library papers are kept as screened context (excluded groups on slide 9)")

# page numbers
for pos,sl in enumerate(prs.slides,1):
    if pos==1: continue
    tb=sl.shapes.add_textbox(IN(0.10),IN(5.26),IN(0.45),IN(0.28))
    tf=tb.text_frame; tf.word_wrap=False
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=str(pos)
    r.font.name=F; r.font.size=Pt(9); r.font.color.rgb=GREY

prs.save(OUTP)
import shutil; shutil.copy(OUTP, os.path.expanduser("~/Downloads/slide_v7.pptx"))
print("saved", OUTP, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
