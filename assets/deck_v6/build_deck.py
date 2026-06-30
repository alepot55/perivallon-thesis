"""Deck v6 builder — WorldView-3 + Pléiades Neo direction.
Didactic, colourful (Fraternali audience). 16:9. Figures from make_figs.py + fix_figs.py.
Content from docs/01_calls/2026-06-30_deck_revision.md; numbers verified from papers/notes.
Run:  python3 assets/deck_v6/build_deck.py  ->  assets/deck_v6/deck_v6.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path("/home/alepot55/Desktop/uni/Tesi")
FIG = ROOT / "assets" / "deck_v6" / "figs"
OUTP = ROOT / "assets" / "deck_v6" / "deck_v6.pptx"

TEAL = RGBColor(0x2E,0x8B,0x8B); TEAL_DK=RGBColor(0x1F,0x6F,0x6F)
ACCENT=RGBColor(0xE0,0x7A,0x3F); INK=RGBColor(0x22,0x22,0x22)
GREY=RGBColor(0x88,0x88,0x88); LGREY=RGBColor(0xF2,0xF2,0xF2)
ASB=RGBColor(0x9B,0x3D,0x6F); GOLD=RGBColor(0xC9,0xA2,0x27); WHITE=RGBColor(0xFF,0xFF,0xFF)
FONT="Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def _tf(tf, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT):
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold
            r.font.italic=italic; r.font.name=FONT

def box(slide,x,y,w,h,fill=None,line=None,line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(line_w)
    return sp

def txt(slide,x,y,w,h,text,size=18,color=INK,bold=False,italic=False,
        align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,line_spacing=1.0):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    lines=text.split("\n") if isinstance(text,str) else text
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=ln; p.alignment=align; p.line_spacing=line_spacing
        p.space_after=Pt(3)
        for r in p.runs:
            r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold
            r.font.italic=italic; r.font.name=FONT
    return tb

def bullets(slide,x,y,w,h,items,size=16,color=INK,gap=5):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        lvl = it[0] if isinstance(it,tuple) else 0
        s   = it[1] if isinstance(it,tuple) else it
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bullet = "    – " if lvl else "•  "
        p.text=bullet+s; p.line_spacing=1.06; p.space_after=Pt(gap)
        for r in p.runs:
            r.font.size=Pt(size-(1 if lvl else 0)); r.font.name=FONT
            r.font.color.rgb=(GREY if lvl else color)
    return tb

def header(slide,title,sub=None):
    box(slide,0,0,SW,Inches(0.12),fill=TEAL)
    txt(slide,Inches(0.55),Inches(0.30),Inches(12.2),Inches(0.9),title,
        size=26,color=TEAL_DK,bold=True)
    if sub:
        txt(slide,Inches(0.55),Inches(1.12),Inches(12.2),Inches(0.5),sub,
            size=15,color=ACCENT,italic=True)

def footer(slide,text):
    txt(slide,Inches(0.55),Inches(7.05),Inches(12.2),Inches(0.35),text,
        size=10,color=GREY,italic=True)

def imgfit(slide,path,x,y,maxw,maxh):
    im=Image.open(path); iw,ih=im.size; ar=iw/ih
    w=maxw; h=int(w/ar)
    if h>maxh: h=maxh; w=int(h*ar)
    px=x+(maxw-w)//2; py=y+(maxh-h)//2
    slide.shapes.add_picture(str(path),px,py,width=w,height=h)

# ── slide builders ────────────────────────────────────────────────────────
def s_title(t,sub,author):
    s=prs.slides.add_slide(BLANK)
    box(s,0,0,SW,SH,fill=WHITE)
    box(s,0,Inches(2.4),SW,Inches(0.10),fill=TEAL)
    txt(s,Inches(0.9),Inches(2.7),Inches(11.5),Inches(2.0),t,size=38,color=TEAL_DK,bold=True)
    txt(s,Inches(0.9),Inches(4.7),Inches(11.5),Inches(0.7),sub,size=20,color=ACCENT,italic=True)
    txt(s,Inches(0.9),Inches(6.4),Inches(11.5),Inches(0.6),author,size=15,color=GREY)
    return s

def s_section(n,t):
    s=prs.slides.add_slide(BLANK)
    box(s,0,0,SW,SH,fill=TEAL)
    txt(s,Inches(0.9),Inches(2.7),Inches(2),Inches(1.5),n,size=80,color=RGBColor(0x7FB,0xCF,0xCF) if False else RGBColor(0xA6,0xCF,0xCF),bold=True)
    txt(s,Inches(0.9),Inches(4.0),Inches(11.5),Inches(1.6),t,size=34,color=WHITE,bold=True)
    return s

def s_text(title,items,sub=None,foot=None,size=17):
    s=prs.slides.add_slide(BLANK); header(s,title,sub)
    bullets(s,Inches(0.7),Inches(1.7),Inches(12.0),Inches(5.0),items,size=size,gap=8)
    if foot: footer(s,foot)
    return s

def s_img(title,img,items=None,sub=None,foot=None,fig_titled=False):
    s=prs.slides.add_slide(BLANK)
    if items:
        header(s,title,sub)
        bullets(s,Inches(0.7),Inches(1.7),Inches(4.3),Inches(5.0),items,size=14,gap=6)
        imgfit(s,FIG/img,Inches(5.2),Inches(1.55),Inches(7.7),Inches(5.2))
    elif fig_titled:
        box(s,0,0,SW,Inches(0.12),fill=TEAL)
        imgfit(s,FIG/img,Inches(0.5),Inches(0.45),Inches(12.3),Inches(6.3))
    else:
        header(s,title,sub)
        imgfit(s,FIG/img,Inches(0.8),Inches(1.55),Inches(11.7),Inches(5.25))
    if foot: footer(s,foot)
    return s

def s_table(title,rows,sub=None,foot=None,colw=None,fs=13):
    s=prs.slides.add_slide(BLANK); header(s,title,sub)
    nr=len(rows); nc=len(rows[0])
    tw=Inches(12.0); th=Inches(0.5*nr if nr*0.5<4.8 else 4.8)
    gx=Inches(0.7); gy=Inches(1.8)
    gt=s.shapes.add_table(nr,nc,gx,gy,tw,th).table
    if colw:
        for j,wv in enumerate(colw): gt.columns[j].width=Inches(wv)
    for i,row in enumerate(rows):
        for j,cell in enumerate(row):
            c=gt.cell(i,j); c.text=str(cell)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.06)
            c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size=Pt(fs+1 if i==0 else fs); r.font.name=FONT
                    r.font.bold=(i==0); r.font.color.rgb=(WHITE if i==0 else INK)
            if i==0: c.fill.solid(); c.fill.fore_color.rgb=TEAL
            else: c.fill.solid(); c.fill.fore_color.rgb=(WHITE if i%2 else RGBColor(0xEF,0xF6,0xF6))
    if foot: footer(s,foot)
    return s

# ════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════
s_title("Multispectral satellite imagery for illegal waste material classification",
        "State of the art and thesis direction — WorldView-3 + Pléiades Neo",
        "Alessandro Potenza · M.Sc. Computer Science & Engineering (AI) · Politecnico di Milano · PERIVALLON")

# SECTION 1
s_section("1","The problem: from detection to risk")
s_text("Illegal waste dumps: a hidden, widespread environmental crime",
    ["Lombardy and Italy face thousands of unauthorised dumping sites — many never reported, only stumbled upon.",
     "ARPA (the regional environmental agency) must find them and decide where to send a limited number of inspectors first.",
     "The operational question is not only “where is the waste?” but “how dangerous is it?” — classify by risk, not merely detect.",
     (1,"Inspector time is the scarce resource; without ranking, agencies inspect blindly.")],
    sub="Define: EO = Earth Observation · MS = multispectral · GSD = ground sampling distance · VNIR/SWIR = visible-near / short-wave infrared",
    foot="Framing: Fraternali et al. 2024 (PoliMi survey, arXiv:2402.09066)")
s_text("Not all waste is equal: risk = hazard × exposure × magnitude",
    ["The danger is material-bound: inert rubble vs plastics vs asbestos-cement are worlds apart.",
     "“Material” carries the hazard; the European Waste Catalogue (EWC) flags hazardous codes with “*” — asbestos = 17 06 05*.",
     "Priority risk = hazard (material) × exposure (who/what is nearby) × magnitude (how much).",
     (1,"So material identification is the missing decision variable, not just site detection.")],
    foot="EWC List of Waste (Dec. 2000/532/EC) · Indice di Degrado d.d.g. 13237/2008 · Fazzo et al. 2023")
s_img("From pixels to an ARPA intervention priority","risk_chain.png",
    foot="Computer vision reframed as decision support · EWC 2000/532/EC · Indice di Degrado d.d.g. 13237/2008", fig_titled=True)
s_text("The research question",
    ["What is the added value of multispectral (MS) data, vs RGB only, for waste MATERIAL classification from satellite imagery?",
     (1,"Measured, not assumed — and explicitly bounded: where does multiband help, and where does it stop?"),
     (1,"Two confirmed very-high-resolution sensors: WorldView-3 (with SWIR) and Pléiades Neo (VNIR-only).")],
    foot="The core scientific bet of the thesis", size=19)

# SECTION 2
s_section("2","Why ordinary RGB is not enough")
s_text("Today’s paradigm: RGB deep learning detects sites, not materials",
    ["Pipeline: VHR RGB tile → CNN/Transformer backbone → EO pretraining → two-step fine-tune → binary waste / no-waste.",
     "Strong within the geographic context it was trained on: Gibellini et al. 2025 — F1 92.0 %, Acc 94.6 % (20 cm, Swin-T + RSP).",
     "But three limits motivate the spectrum:",
     (1,"classifies presence, not material;"),
     (1,"tied to colour only (3 broad bands);"),
     (1,"generalizzazione collapses cross-region (−5.1 % F1: GR 85.5 / SE 83.8 / RO 91.5) and cross-sensor.")],
    foot="Gibellini et al. 2025 (Waste Mgmt Bull., arXiv:2502.06607) · Fraternali et al. 2024")
s_text("A satellite pixel is a spectrum, not just a colour",
    ["Each pixel is a vector of reflectance values — the material’s spectral signature.",
     "RGB: 3 broad visible bands (colour only).",
     "MS (multispectral): 4–15 chosen bands. WorldView-3 records 8 VNIR + 8 SWIR; Pléiades Neo 6 VNIR.",
     "HSI (hyperspectral): hundreds of contiguous narrow bands (e.g. EnMAP 230, splib07a 2,151 channels @ 350–2500 nm).",
     (1,"This vector is the raw input from which any classifier reasons.")],
    foot="Spectrum schema after USGS splib07a (Kokaly et al., 2017)")
s_img("Every material has a spectral fingerprint","fingerprint_4_materials.png",
    foot="Reflectance curves diagnose chemistry, not just colour — USGS splib07a (Kokaly et al., 2017)")
s_img("RGB fails in two distinct ways","rgb_fails_two_panels.png",
    foot="Iso-chromaticity + sub-pixel mixing — Tasseron 2021 · Aguilar/Uhrin 2025 · USGS splib07a")
s_img("Where the diagnostic information lives","where_info_lives.png",
    foot="The chemistry-diagnostic clues sit in NIR and especially SWIR — USGS splib07a", fig_titled=True)

# SECTION 3
s_section("3","The sensors — what exists today")
s_img("The sensor trade-off: spatial × spectral × revisit","sensor_radar.png",
    foot="Three axes, one photon budget — no satellite maximises all three", fig_titled=True)
s_table("Multispectral at low / medium resolution (10–30 m): too coarse for material",
    [["Sensor","Bands","GSD","SWIR","Verdict for material"],
     ["Sentinel-2","13","10–20 m","20 m","detection only; pixels too mixed"],
     ["PRISMA / EnMAP","200+ (HSI)","30 m","yes","resolves 2.3 µm triplet but 30 m"],
     ["Landsat 8/9","11","30 m","yes","regional context only"]],
    sub="Low resolution answers “where”, not “what”",
    foot="Kikaki 2022 · Magyar 2023 · Shepherd 2025 · Sun 2023",
    colw=[2.6,1.6,1.7,1.4,4.7], fs=14)
s_text("Low / medium resolution in practice",
    ["MARIDA (Kikaki 2022): marine debris confused with natural organic matter where pixels mix.",
     "Global dumpsites (Sun 2023, Sentinel-2): 763 dumpsite LOCATIONS across 28 cities (sensitivity 98 %, precision 70 %) — but material identity stays out of reach; mapping in 6 days vs 6 months manual.",
     "Tisza (Magyar 2023): change detection works; per-tile post-processing still needed.",
     (1,"Takeaway: low-res localises sites; it cannot tell what they are made of. → we need VHR.")],
    foot="Kikaki 2022 PLOS ONE · Sun 2023 Nat. Commun. 14 · Magyar 2023")
s_img("The chosen data: WorldView-3 + Pléiades Neo","spec_cards.png",
    foot="Both sub-metre VHR · free via ESA Third-Party Missions (~9-week proposal lead)", fig_titled=True)
s_img("Honest caveat: resolution is split between texture and chemistry","swir8_bottleneck.png",
    items=["Spectra sampled coarsely: VNIR ~1.24 m, SWIR ~3.7 m (~14 m² pixels — no pure pixels over a dump).",
     "Pan gives fine texture ~0.3 m but no material info alone.",
     "A small dump is easier to SEE (texture) than to chemically IDENTIFY (spectra).",
     "SWIR-8 bottleneck: asbestos 2.32 + concrete 2.34 + plastic 2.31 crowd one band.",
     "Honest: at 3.7 m + atmosphere, R3 may ≈ R2 — itself a valid finding."],
    foot="USGS splib07a · WV-3 SWIR band centres (Maxar)")

# SECTION 4
s_section("4","The added value is real — precedents")
s_img("Spectral added value, measured: Aguilar 2021 on WorldView-3","aguilar_bars.png",
    items=["WV-3 greenhouse-plastic ablation, OBIA + Decision Tree, 10-fold CV.",
     "Overall Accuracy:",
     (1,"VNIR only → 90.85 %"),
     (1,"SWIR only → 96.79 %"),
     (1,"All 16 bands → 97.38 %"),
     "SWIR carries the single largest jump (+6.5 pp over VNIR).",
     "Most-cited direct benchmark for waste-like material classification."],
    foot="Aguilar, Jiménez-Lao & Aguilar 2021, Remote Sensing 13(11):2133 — κ 0.812/0.932/0.944")
s_img("More bands ≠ more — well-chosen few suffice","bands_plateau.png",
    items=["Aguilar 2021 (WV-3 plastic): VNIR 90.85 → All 97.38 — SWIR carries the jump.",
     "Vitek 2025 (C&D waste): RGB + 2 narrowbands ≈ HSI 768 bands.",
     "Zhou 2021 (WV-3 plastic): 8 SWIR bands separate 3 polymer clusters.",
     "What matters is WHICH bands, not how many.",
     "WV-3 gives us SWIR; Pléiades Neo (VNIR-only) measures how far we get without it."],
    foot="Vitek et al. 2025 (CTU Prague) · Aguilar et al. 2021 · Zhou et al. 2021")
s_text("Asbestos precedents — direct comparison points",
    ["Saba 2026 (WV-3 VNIR, 8 bands): 32 classifiers tested; Fine-KNN Macro-F1 97.6 % — strong VNIR-only upper bound to beat. [paywalled, single-source]",
     "Bonifazi 2026 (WV-3, 16 bands, MLC): asbestos-cement roofs F1 0.87 (P 0.81 / R 0.92); 25,319 buildings processed; multi-temporal monitoring.",
     "Shepherd 2025 (EnMAP, 230 bands @ 30 m): ACE classifier 91.4 % OA / κ 0.87 / F1 91.2 %; 86 % field match — but urban mixed pixels (rubble, paint) mimic asbestos.",
     (1,"Together: VNIR can detect AC; SWIR + hyperspectral add robustness on clean/ambiguous roofs.")],
    foot="Saba 2026 J. Hazard. Mater. · Bonifazi 2026 Geomatics 6(3):41 · Shepherd 2025 Sci. Rep.")
s_text("Waste-dump precedents at VHR",
    ["CascadeDumpNet (Zhang & Ma 2024, Pléiades 0.5 m): open-dumpsite detection, 84.6 % mAP; transfers Shenzhen → Shanghai/Guangzhou.",
     "AerialWaste (Torres 2023): 11,700+ RGB tiles, 22 material categories — but only 11 images contain presumed asbestos; RGB, object-level.",
     "Aguilar 2025 (WV-3 SWIR macroplastics): precision 92.5 %, lab-to-satellite r = 0.95; targets 80–150 m² aggregations.",
     (1,"VHR detection is solved on RGB; material/hazard identification is the open frontier.")],
    foot="Zhang & Ma 2024 RSE · Torres 2023 Sci. Data · Aguilar 2025 Env. Monit. Assess.")
s_img("Where each hazard becomes separable: feature → band","band_material_map.png",
    foot="Pléiades Neo reaches up to NIR; WorldView-3 adds the SWIR chemistry — USGS splib07a · Aguilar · Cilia 2015", fig_titled=True)

# SECTION 5
s_section("5","Method — a fair backbone")
s_text("Foundation models for Earth Observation: a new lever",
    ["2024–2025: a wave of pretrained EO backbones (Prithvi, SpectralGPT, SoftCon, AnySat, DOFA…).",
     "Most are pretrained at 10–30 m on Sentinel-2 / HLS — naïve transfer to ~1.2 m VHR is not guaranteed.",
     "Two routes forward:",
     (1,"sensor-agnostic models (DOFA, AnySat) accept arbitrary band sets;"),
     (1,"parameter-efficient adapters (DEFLECT) freeze the backbone, add < 1 % params."),
     "Open question: does 10–30 m pretraining transfer to VHR? — empirical."],
    foot="FM-RS survey 2024 · Xiong 2024 · Astruc 2024 (AnySat) · Thoreau 2025 (DEFLECT) · Szwarcman 2024 (Prithvi-EO-2.0)")
s_img("DOFA: a band-agnostic backbone makes the comparison fair","dofa_schematic.png",
    foot="Xiong et al. 2024 (arXiv:2403.15356) — pretrained on 5 modalities (8M samples), handles up to 202 bands", fig_titled=True)

# SECTION 6
s_section("6","The gap & the proposed direction")
s_text("Gaps in the literature — and the test that matters",
    ["No public VHR dataset pairs waste sites with MATERIAL labels; no peer-reviewed Pléiades-Neo material benchmark exists.",
     "RGB pipelines conflate iso-colour, chemically distinct materials.",
     "No pipeline goes from spectra to a regulatory hazard class (material → EWC → risk tier).",
     "The object-vs-material split is unaddressed; generalizzazione (cross-region / sensor / time) is rarely shown.",
     "The test that matters: how much does multiband actually beat RGB for material — and where does it stop helping?"],
    foot="Gap list distilled from Fraternali et al. 2024 (arXiv:2402.09066)")
s_img("The experiment: a 3-axis band ablation","ablation_cube.png",
    foot="Per-pixel AND full-CNN: B−A = pure chemistry gain, D−C = total MS gain · design in loop_prof_sota/04", fig_titled=True)
s_img("Generalization as a first-class axis","generalization_2d.png",
    foot="Report ID−OOD ΔF1 per cell — pattern EXPECTED, to be measured · cross-region −5.1 %: Gibellini 2025", fig_titled=True)
s_text("From material to risk tier",
    ["Map predicted material → EWC hazard code → Indice di Degrado / PRAL logic × exposure × magnitude → ARPA priority tier.",
     "Output = a ranked intervention list; every step transparent and auditable.",
     (1,"Indice di Degrado (ID ≤25 / 26–44 / ≥45 → action windows) is NOT in the public WFS → estimated remotely (SWIR Mg-OH depth + VNIR weathering) = a thesis contribution."),
     (1,"Exposure grounded in Fazzo et al. (residential-proximity mesothelioma risk); magnitude from roof/site area.")],
    foot="EWC 17 06 05* · Indice di Degrado d.d.g. 13237/2008 · Fazzo et al. 2023")
s_img("The asbestos pilot: the immediately-feasible demonstrator","pilot_workflow.png",
    foot="Public pixel-accurate labels + textbook SWIR diagnostic (Mg-OH 2.30–2.33 µm) — de-risks the whole pipeline", fig_titled=True)
s_text("Why this matters — and what is genuinely novel",
    ["First to MEASURE RGB vs VNIR vs SWIR added value for waste material at VHR.",
     "First WorldView-3 vs Pléiades Neo head-to-head for this task (SWIR vs VNIR-only).",
     "First to connect spectra → EWC hazard → ARPA risk tier — confidence-aware (reports where MS helps and where it doesn’t, incl. generalizzazione).",
     (1,"Payoff: a transferable, free-data decision-support tool for environmental agencies.")],
    foot="Contributions map one-to-one onto the literature gaps")
s_text("What to confirm with Thomas — decision checklist",
    ["Data: trigger the ESA TPM proposal now (~9-week lead)? AOI / scene selection over Lombardia?",
     "Risk-chain depth: full Indice di Degrado / PRAL formalization, or stop at material + EWC?",
     "SuperDove: kept as a free wide-area screening layer, or dropped?",
     "Ablation depth: 3 axes as designed, or trim for time?",
     "Backbone: DOFA primary + Swin-RSP RGB reference — agreed?",
     "AerialWaste coordinate bridge with ARPA (gates the general-waste MS ablation)."],
    foot="Open dependencies flagged in loop_prof_sota/00_LOOP_LOG.md")

# APPENDIX
s_section("A","Appendix — backup")
s_table("Sensor comparison (full)",
    [["Sensor","VNIR","SWIR","Pan","Access","Role"],
     ["WorldView-3","8 @ 1.24 m","8 @ 3.7 m","0.31 m","commercial / TPM","full ladder incl. R3"],
     ["Pléiades Neo","6 @ 1.2 m","none","0.30 m","commercial / TPM","VNIR-only cross-sensor"],
     ["SuperDove","8 @ 3 m","none","—","free, near-daily","wide-area screening"],
     ["Sentinel-2","10–20 m","20 m","—","free","too coarse for material"],
     ["EnMAP / PRISMA","HSI ~6.5 nm","yes","—","free (proposal)","30 m — context"]],
    foot="Specs: Maxar · Airbus · Planet · ESA · ASI · DLR",
    colw=[2.4,1.9,1.5,1.2,2.4,2.6], fs=12)
s_text("The 13 target materials: shape vs chemistry",
    ["Shape-identifiable (8 — RGB usually enough): vehicles, tanks, containers, rubble, bulky, wood, tyres, big-bags.",
     "Chemistry-bound (5 — RGB blind, MS lives here): asbestos, plastic-type, foundry slag, sludge, scrap-metal composition.",
     (1,"Asbestos + plastics are ideal test beds — shape cue AND chemistry cue coexist.")],
    foot="Taxonomy: loop_prof_sota/09_master_cheatsheet.md")
s_table("Diagnostic feature → band (memorise)",
    [["Material","Feature (µm)","WV-3 band","Pléiades Neo?"],
     ["Asbestos (Mg-OH)","2.30–2.33","SWIR-7 (2.33)","✗"],
     ["Plastics (C-H)","1.215 / 1.73 / 2.31","S1, S4, S8","✗"],
     ["Concrete / C&D","2.34 / 2.20","S8 / S6","✗"],
     ["Slag / rust (Fe-oxide)","0.87–0.95","NIR1 / NIR2","weak"],
     ["AC weathering","0.68 + 0.74","Red + RedEdge","✓"]],
    foot="USGS splib07a · Aguilar 2021/2025 · Cilia 2015",
    colw=[3.0,2.8,2.6,2.4], fs=13)
s_text("Honesty guardrails & claims to handle with care",
    ["Multiband value is measured, not assumed — and bounded (SWIR earns its keep on chemistry; detection may need little beyond RGB).",
     "Spectra at 1.24 m / 3.7 m vs pan 0.3 m = texture, not chemistry; SWIR carries chemistry but Pléiades Neo has none.",
     "Unmixing is out of scope (it sets the ceiling). “generalizzazione”, not “OOD”.",
     (1,"Downgraded / do-not-quote: Aguilar “+14 % κ” (use OA only); Saba “HS 97.3 vs MS 74.4” (paywalled); chrysotile “~5 nm” (it is tens of nm).")],
    foot="loop_prof_sota/09_master_cheatsheet.md")
s_text("Key references",
    ["Gibellini et al. 2025 — DL waste pipeline (Swin-T+RSP baseline). Waste Mgmt Bull.",
     "Fraternali et al. 2024 — waste RS survey / gap analysis. arXiv:2402.09066.",
     "Aguilar et al. 2021 — WV-3 VNIR+SWIR ablation. Remote Sensing 13(11):2133.",
     "Xiong et al. 2024 — DOFA band-agnostic foundation model. arXiv:2403.15356.",
     "Bonifazi 2026 (Geomatics) · Saba 2026 (J. Hazard. Mater.) · Shepherd 2025 (Sci. Rep.) — asbestos.",
     "Zhang & Ma 2024 — CascadeDumpNet (Pléiades). RSE. · Kokaly et al. 2017 — USGS splib07a.",
     "Full annotated bibliography: docs/02_research/loop_prof_sota/11_references.md + references.bib"],
    foot="Politecnico di Milano · PERIVALLON Horizon Europe (Grant 101073952)", size=15)

prs.save(str(OUTP))
print("saved", OUTP, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
