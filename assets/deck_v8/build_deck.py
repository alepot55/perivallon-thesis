"""Deck v8. Restructured after Thomas's feedback on v7 (2026-07-09):
linear section order, literature search told as method -> results -> selection,
SOTA analysis in a single block ending on Alari 2024 (the anchor reference,
now two slides), proposal material all after the gaps. 28 slides; all 47
library papers cited (References split by theme over three slides)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE=os.path.dirname(os.path.abspath(__file__))
FIG=os.path.join(HERE,"figs")
OUTP=os.path.join(HERE,"deck_v8.pptx")
INK=RGBColor(0x1A,0x1A,0x1A); GREY=RGBColor(0x66,0x66,0x66)
HDR=RGBColor(0xDD,0xDD,0xDD); ALT=RGBColor(0xF5,0xF5,0xF5); WHITE=RGBColor(0xFF,0xFF,0xFF)
F="Calibri"
prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
BLANK=prs.slide_layouts[6]
def IN(v): return Inches(v)
def tf_fill(tf,items,wrap=True):
    tf.clear(); tf.word_wrap=wrap
    for k,it in enumerate(items):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(it.get("sa",6)); p.line_spacing=it.get("ls",1.12)
        r=p.add_run(); r.text=it["t"]
        r.font.name=F; r.font.size=Pt(it.get("sz",14)); r.font.bold=it.get("b",False)
        r.font.italic=it.get("i",False); r.font.color.rgb=it.get("c",INK)
def slide(sec=None):
    s=prs.slides.add_slide(BLANK)
    if sec:
        tb=s.shapes.add_textbox(IN(7.35),IN(0.10),IN(2.45),IN(0.28))
        tf=tb.text_frame; tf.word_wrap=False
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=sec
        r.font.name=F; r.font.size=Pt(9); r.font.italic=True; r.font.color.rgb=GREY
    return s
def title(s,t):
    tb=s.shapes.add_textbox(IN(0.45),IN(0.28),IN(9.10),IN(0.6))
    tf_fill(tb.text_frame,[{"t":t,"b":True,"sz":22}])
def body(s,items,top=1.20,left=0.45,w=9.10,h=3.85,mid=True):
    tb=s.shapes.add_textbox(IN(left),IN(top),IN(w),IN(h)); tf_fill(tb.text_frame,items)
    if mid: tb.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return tb
def foot(s,t):
    tb=s.shapes.add_textbox(IN(0.45),IN(5.26),IN(9.10),IN(0.26))
    tf_fill(tb.text_frame,[{"t":t,"sz":8,"c":GREY,"sa":0}])
def img(s,path,left,top,w,h):
    from PIL import Image
    iw,ih=Image.open(path).size; ar=iw/ih
    W=IN(w); H=int(W/ar)
    if H>IN(h): H=IN(h); W=int(H*ar)
    s.shapes.add_picture(path,IN(left)+(IN(w)-W)//2,IN(top)+(IN(h)-H)//2,width=W,height=H)
def table(s,data,left=0.45,top=1.25,w=9.10,colw=None,fs=9.5,rh=0.46,hh=0.38,vc=False):
    nr=len(data); nc=len(data[0]); th=hh+(nr-1)*rh
    if vc: top=1.12+max(0,(4.0-th)/2.0)
    gt=s.shapes.add_table(nr,nc,IN(left),IN(top),IN(w),IN(th)).table
    if colw:
        for j,cw in enumerate(colw): gt.columns[j].width=IN(cw)
    gt.rows[0].height=IN(hh)
    for i in range(1,nr): gt.rows[i].height=IN(rh)
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.text=str(val); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=IN(0.05); c.margin_right=IN(0.04); c.margin_top=IN(0.015); c.margin_bottom=IN(0.015)
            c.fill.solid(); c.fill.fore_color.rgb=HDR if i==0 else (ALT if i%2==0 else WHITE)
            for p in c.text_frame.paragraphs:
                p.line_spacing=1.0
                for r in p.runs:
                    r.font.name=F; r.font.size=Pt(fs+0.5 if i==0 else fs)
                    r.font.bold=(i==0); r.font.color.rgb=INK
    return gt

# ── Part I: problem and task ──────────────────────────────────────────────

# 1 ── title
s=slide()
tb=s.shapes.add_textbox(IN(0.7),IN(1.7),IN(8.6),IN(1.4))
tf_fill(tb.text_frame,[{"t":"Classification of waste materials in very-high-resolution satellite imagery","b":True,"sz":30}])
tb2=s.shapes.add_textbox(IN(0.7),IN(3.15),IN(8.6),IN(0.5))
tf_fill(tb2.text_frame,[{"t":"State of the art and thesis proposal","sz":17,"i":True,"c":GREY}])
tb3=s.shapes.add_textbox(IN(0.7),IN(4.5),IN(8.6),IN(0.4))
tf_fill(tb3.text_frame,[{"t":"Alessandro Potenza  ·  M.Sc. Computer Science and Engineering, AI  ·  PERIVALLON","sz":12,"c":GREY}])

# 2 ── outline
s=slide(); title(s,"Outline")
body(s,[
 {"t":"1.  Problem and task: what has to be recognised, in which images, over which materials (slides 3-7).","sz":14,"sa":12},
 {"t":"2.  Literature search: how the search was run, what it returned, what was kept and why (slides 8-10).","sz":14,"sa":12},
 {"t":"3.  State of the art: the relevant works, from site-level detection down to the direct predecessor, and what is missing (slides 11-21).","sz":14,"sa":12},
 {"t":"4.  Proposed work: approach, available imagery, asbestos pilot, evaluation (slides 22-25).","sz":14,"sa":12},
])
foot(s,"Each part builds on the previous one; the proposal is derived from the gaps at the end of part 3.")

# 3 ── context
s=slide("Problem and task"); title(s,"Context")
body(s,[
 {"t":"Illegal waste dumping is an environmental crime with direct public-health consequences. Agencies (ARPA) have limited inspection capacity, and priority depends on what is dumped: rubble, plastics and asbestos-cement imply very different hazards.","sz":12.5},
 {"t":"Detecting dump sites from images is mature. Recognising the material is not: the reference survey (50 works, 1987-2023, almost all RGB) marks it as an open problem.","sz":12.5},
 {"t":"One thesis in the group has addressed material classification directly (Alari 2024). This work starts from there.","sz":12.5},
],top=1.00,h=1.85,mid=False)
img(s,os.path.join(FIG,"torres_examples.png"),0.7,2.95,8.6,2.20)
foot(s,"Images: AerialWaste positive examples, Torres and Fraternali 2023 (CC BY 4.0)  ·  survey: Fraternali et al. 2024")

# 4 ── task definition
s=slide("Problem and task"); title(s,"Task definition")
body(s,[
 {"t":"Objective: given a VHR satellite image of a suspected dump site, output the set of waste materials present.","b":True,"sz":13},
 {"t":"Task: multi-label classification at image level.","sz":13},
 {"t":"Input: VHR optical imagery, GSD 0.2 to 1.3 m. Aerial RGB for the baseline (AerialWaste); satellite VNIR for the multispectral arm, up to 8 bands. SWIR excluded: not in the planned acquisitions, and its 3.7 m GSD does not match the task. Sentinel-2 excluded: too coarse (10-20 m).","sz":13},
 {"t":"Technique: classification, not object detection or segmentation. Annotations are image-level (Alari 2024: 11,477 multi-label); waste piles have no stable shape for boxes; no segmentation masks exist for these datasets.","sz":13},
 {"t":"Research question: does VNIR input improve waste-material classification over RGB, and for which materials?","b":True,"sz":13},
])
foot(s,"Annotations: Alari 2024  ·  imagery: WorldView-3, Pléiades Neo (VNIR + pan)")

# 5 ── task scheme
s=slide("Problem and task"); title(s,"Task scheme")
img(s,os.path.join(FIG,"task_io.png"),0.7,1.35,8.6,3.1)
body(s,[{"t":"One image in, a set of material labels out. The same scheme runs with RGB or VNIR input; only the first layer of the network changes.","sz":12}],top=4.55,h=0.6,mid=False)
foot(s,"Tile: Alari 2024. The multi-label formulation follows the group predecessor.")

# 6 ── materials taxonomy
s=slide("Problem and task"); title(s,"The material taxonomy")
img(s,os.path.join(FIG,"alari_row1.png"),0.85,1.10,8.3,1.75)
img(s,os.path.join(FIG,"alari_row2.png"),0.85,2.95,8.3,1.75)
body(s,[{"t":"Ten-category grouping used in the group predecessor; the full annotation set counts 13 categories (adds foundry waste, sludge, tanks as separate classes).","sz":11.5,"c":GREY}],top=4.72,h=0.5,mid=False)
foot(s,"Samples and taxonomy: Alari 2024 (politesi 10589/230633), built on AerialWaste annotations")

# 7 ── materials decision
s=slide("Problem and task"); title(s,"Materials: which are considered and why")
table(s,[
 ["Material","Target","Spectral analysis","Reason"],
 ["Rubble, plastic, wood, tires","yes","yes","frequent and confusable in RGB"],
 ["Asbestos-cement","yes","yes, dedicated pilot","public regional ground truth (WFS)"],
 ["Vehicles, tanks, containers, scrap, bulky, big bags","yes","no","shape-based; RGB sufficient in literature"],
 ["Sludge, foundry waste","yes","no","few labels; visually ambiguous at this GSD"],
],left=0.45,top=1.45,w=9.10,colw=[3.4,0.95,1.85,2.9],fs=9.5,rh=0.52,hh=0.40)
body(s,[{"t":"All 13 categories remain classification targets, for continuity with the group dataset. The RGB-versus-VNIR analysis concentrates on the subset where colour is ambiguous and labels are sufficient.","sz":12}],top=4.15,h=0.85,mid=False)
foot(s,"Decision criteria: visual ambiguity in RGB, label availability, hazard relevance")

# ── Part II: literature search ────────────────────────────────────────────

# 8 ── how the search was run
s=slide("Literature search"); title(s,"Literature search: how it was run")
body(s,[
 {"t":"Scripted queries on the Scopus API (waste detection in remote sensing; asbestos roof mapping), deduplication, manual screening with explicit criteria, snowballing from the Fraternali 2024 survey and the Alari 2024 references.","sz":12.5},
],top=1.05,h=1.0,mid=False)
img(s,os.path.join(FIG,"search_flow.png"),0.55,2.15,8.9,2.55)
foot(s,"Search artifacts: papers/literature_search (queries, raw and deduplicated records)  ·  script: papers/scripts/scopus_search.py")

# 9 ── what the search returned
s=slide("Literature search"); title(s,"What the search returned")
table(s,[
 ["Topic","Works at task scale","Notes"],
 ["Waste sites, no material distinction","Torres 2023, Gibellini 2025, Sun 2023, CascadeDumpNet 2024, CWLD 2024, Disaitek 2024","mature line; answers where, not what"],
 ["Multi-material waste classification","Alari 2024 (group thesis)","the only direct precedent"],
 ["Asbestos-cement","Saba 2026, Bonifazi 2026, Abbasi 2024, Cilia 2015","roofs only, never inside a waste taxonomy"],
 ["Tanks","Ramachandran 2024, YOLOv7-OT 2024","mature object detection"],
 ["Vehicles, containers","vehicle DA 2020, truck-and-container 2025","VHR object detection"],
 ["Scrap","ELV Hybrid-YOLOv5 2025","close-range infrared, not satellite"],
 ["Rubble / inert","CWLD 2024; debris volume UAV 2022","C&D segmentation; UAV photogrammetry"],
 ["Plastic, wood, tires, bulky, big bags, sludge, foundry","none found in scope","only as classes in AerialWaste / Alari"],
],left=0.45,top=1.14,w=9.10,colw=[2.55,3.75,2.80],fs=8.5,rh=0.44,hh=0.36)
foot(s,"Scope: terrestrial, GSD compatible with the task. Site detection is well covered; material-level literature is thin.")

# 10 ── kept vs excluded
s=slide("Literature search"); title(s,"What was kept, what was excluded")
table(s,[
 ["Excluded group","Examples","Reason"],
 ["Sentinel-2 and marine debris","MARIDA 2022, Tisza 2023","10-20 m: pixels too mixed; different domain"],
 ["Spaceborne hyperspectral","Shepherd 2025 (EnMAP), EMIT 2025","material evidence, but 30-60 m GSD"],
 ["SWIR-based VHR","Aguilar 2021, Guo and Li 2020, Zhou 2021, Aguilar 2025","SWIR not in our acquisitions; 3.7 m GSD"],
 ["Laboratory / spectral libraries","Vitek 2025, SpectralWaste 2024, Knaeps 2020","not Earth observation; kept as band evidence"],
 ["EO foundation models","DOFA 2024, AnySat 2025, Prithvi-EO-2 2024, SpectralGPT 2024, +6","pretrained at 10-30 m; sub-metre transfer unproven"],
],left=0.45,top=1.30,w=9.10,colw=[2.55,3.20,3.35],fs=9,rh=0.50,hh=0.38)
body(s,[{"t":"Kept: works on terrestrial waste or roof materials at task-compatible GSD, plus the reference spectral library. The excluded papers stay in the annotated library as context.","sz":11.5}],top=4.35,h=0.7,mid=False)
foot(s,"Full accounting in papers/INDEX.md (47 papers, each with a structured note)")

# ── Part III: state of the art ────────────────────────────────────────────

# 11 ── site-level table
s=slide("State of the art"); title(s,"Related work: site-level waste detection")
table(s,[
 ["Work (year)","Input / GSD","Task","Method","Result"],
 ["Gibellini 2025","aerial RGB, 20 cm","classification","Swin-T, RSP pretraining","F1 92.0; cross-region -5.1"],
 ["AW ensemble 2025 *","aerial RGB","classification","CNN + transformer ensemble","binary F1 92.4"],
 ["Disaitek 2024","Pléiades Neo, 30 cm","detection","operational service","~95% on sites >2 m2 (vendor)"],
 ["CascadeDumpNet 2024","Pléiades, 50 cm","object detection","two-stage CNN + AutoML","mAP 84.6, cross-city transfer"],
 ["Sun 2023","VHR RGB, 0.3-1 m","object detection","BCA-Net (Faster R-CNN)","~2,500 dumpsites, 28 cities"],
 ["CWLD 2024","GF-2 + GE, 0.5-0.8 m","segmentation","DeepLabV3+ variant","F1 88.9, construction waste"],
 ["AerialWaste, Torres 2023","aerial RGB, 20-50 cm","dataset + cls.","ResNet-FPN baseline","F1 80.7; 22 material tags"],
],colw=[2.05,1.85,1.45,1.95,1.80],fs=8.5,rh=0.42,hh=0.36,vc=True)
foot(s,"All RGB: they answer where a site is, not what it contains. * preprint; Disaitek is a vendor figure.")

# 12 ── deep-dive AerialWaste
s=slide("State of the art"); title(s,"AerialWaste: the reference dataset")
body(s,[
 {"t":"10,434 locations from ARPA Lombardia records, 487 municipalities; three sources: AGEA orthophotos (20 cm), WorldView-3 (30 cm), Google Earth (50 cm).","sz":11.5},
 {"t":"Binary site labels for detection, plus 22 material tags (type of visible object + storage mode) annotated by experts.","sz":11.5},
 {"t":"Material tags cover about 72% of positives, but the dataset was published for detection: material tags are metadata, not a benchmark.","sz":11.5},
 {"t":"This is the label base both Gibellini 2025 and Alari 2024 build on.","sz":11.5},
],top=1.15,left=0.45,w=4.10,h=3.90)
img(s,os.path.join(FIG,"aerialwaste_grid.png"),4.75,1.10,4.85,4.00)
foot(s,"Image: annotated-label examples, Torres and Fraternali 2023, Scientific Data (CC BY 4.0)")

# 13 ── deep-dive Gibellini
s=slide("State of the art"); title(s,"Gibellini 2025: the site-level baseline")
body(s,[
 {"t":"Binary waste / no-waste classification on AerialWaste; Swin-T with remote-sensing pretraining, two-step fine-tuning.","sz":11.5},
 {"t":"F1 92.02 in domain. Cross-region generalisation drops 5.1 points on average (Greece 85.4, Serbia 83.8, Romania 91.5).","sz":11.5},
 {"t":"Saliency maps confirm the model looks at waste heaps, but the output is presence only: no material information.","sz":11.5},
 {"t":"This architecture and training recipe is the starting point of the proposal.","sz":11.5},
],top=1.15,left=0.45,w=4.10,h=3.90)
img(s,os.path.join(FIG,"gibellini_cams.png"),4.75,1.30,4.85,3.55)
foot(s,"Image: true positives with saliency overlays, Gibellini et al. 2025 (Greek generalisation set)")

# 14 ── object-level table
s=slide("State of the art"); title(s,"Related work: objects and close range")
table(s,[
 ["Work (year)","Input","Task","Result"],
 ["Ramachandran 2024","VHR satellite","object detection","tanks, P 0.96 / R 0.97, >169k mapped"],
 ["YOLOv7-OT 2024 *","VHR satellite","object detection","tanks, precision 95.9"],
 ["Truck-and-container 2025 *","VHR satellite","object detection","container size and status"],
 ["Vehicle DA 2020 *","VHR 30-50 cm","object detection","+10% with domain adaptation"],
 ["ELV Hybrid-YOLOv5 2025","close-range IR","object detection","scrap metals, mAP 84.2"],
 ["UAV solid waste 2024 *","UAV RGB","segmentation","OA >94, generic waste piles"],
 ["C&D debris UAV 2022 *","UAV + photogrammetry","segmentation","IoU 0.90 + volume estimation"],
 ["fCLIPSeg 2025","aerial RGB","segmentation","debris, Dice 0.70, event-agnostic"],
 ["Plastic UAV-SWIR 2026 *","UAV hyperspectral SWIR","segmentation","plastic waste, cross-domain"],
 ["Plastic UAV-IoT 2025 *","UAV RGB","object detection","plastic waste, edge deployment"],
],colw=[2.30,1.80,1.75,3.25],fs=8.5,rh=0.38,hh=0.34,vc=True)
foot(s,"Shape-defined objects are mature; material composition still needs spectra or close range (UAV line surveyed in Drones 2025). * pending full-text verification")

# 15 ── deep-dive asbestos
s=slide("State of the art"); title(s,"The asbestos line, in detail")
body(s,[
 {"t":"Saba 2026: WorldView-3, VNIR only, 32 classifiers compared; best Macro-F1 97.6 per-pixel. Red Edge and NIR carry the discrimination.","sz":11},
 {"t":"Bonifazi 2026: WorldView-3 multi-temporal workflow; building-level decisions from pixel classification; tracks roof removals between acquisitions.","sz":11},
 {"t":"Abbasi 2024: aerial RGB with temporal features, OA ~96: shape and time can substitute spectra at fine GSD.","sz":11},
 {"t":"Cilia 2015: airborne hyperspectral, PA 89 / UA 86, plus a weathering index from red/NIR: the degradation angle.","sz":11},
 {"t":"Shepherd 2025: spaceborne hyperspectral (EnMAP, 30 m), field-calibrated library, 86% match on exhaustive ground truth: the signature survives from orbit; the limit is resolution.","sz":11},
 {"t":"Asbestos slate on drone RGB (2023): partially recognisable by shape and texture alone.","sz":11},
 {"t":"Why it matters here: public ground truth (Lombardy WFS, 10,903 roofs) and VNIR evidence make asbestos the natural pilot.","b":True,"sz":11},
],top=1.05,left=0.45,w=4.10,h=4.05)
img(s,os.path.join(FIG,"bonifazi_example.png"),4.85,1.20,4.65,3.85)
foot(s,"Image: building-level asbestos classification on WorldView-3, Bonifazi et al. 2026, Geomatics (CC BY 4.0)")

# 16 ── material-level table
s=slide("State of the art"); title(s,"Related work: material-level classification")
table(s,[
 ["Work (year)","Input / GSD","Task","Result"],
 ["Alari 2024 (PoliMi)","satellite RGB","multi-label cls.","wF1 69.2 (5 cat.), 59.4 (10 cat.)"],
 ["Saba 2026 *","WV-3 VNIR, 1.24 m","pixel cls.","asbestos, Macro-F1 97.6"],
 ["Bonifazi 2026","WV-3 VNIR+SWIR","pixel cls.","asbestos roofs, removal tracking"],
 ["Abbasi 2024 *","aerial RGB","OBIA cls.","asbestos, OA ~96 (shape + time)"],
 ["Cilia 2015","airborne HSI, 3 m","pixel cls.","asbestos, PA 89 / UA 86"],
],left=0.45,top=1.35,w=9.10,colw=[2.35,2.20,1.85,2.70],fs=9.5,rh=0.50,hh=0.40)
body(s,[{"t":"One multi-material predecessor; everything else addresses a single material, asbestos, on roofs. No work measures RGB versus VNIR on waste materials. The predecessor is examined next.","sz":12}],top=4.30,h=0.75,mid=False)
foot(s,"* paywalled or preprint, reported as such")

# 17 ── deep-dive Alari (1): framing
s=slide("State of the art"); title(s,"Alari 2024: the direct predecessor")
body(s,[
 {"t":"The anchor reference of this thesis: the only work that frames waste-material recognition as multi-label classification, done in the same group, on the same imagery base.","b":True,"sz":12},
 {"t":"Dataset: 11,477 multi-label annotations over 13 categories; 3,190 positive and 7,190 negative images, built on AerialWaste imagery and ARPA records.","sz":12},
 {"t":"Models: ResNet-50 and Swin backbones with FPN; three classification-head designs; weighted binary cross-entropy for label imbalance; different pretraining sources compared.","sz":12},
 {"t":"The taxonomy on slide 6 and the target materials on slide 7 come directly from this work.","sz":12},
 {"t":"Headline results: weighted F1 69.21 on five categories, 59.42 on ten. Input is RGB only; the spectral dimension is untouched.","sz":12},
])
foot(s,"Alari 2024, M.Sc. thesis, PoliMi (advisor Fraternali, co-advisor Gibellini), politesi 10589/230633")

# 18 ── deep-dive Alari (2): per-category results
s=slide("State of the art"); title(s,"Alari 2024: results by category")
body(s,[
 {"t":"Growing the taxonomy from 5 to 10 categories costs 9.8 weighted-F1 points (69.2 to 59.4).","sz":11.5},
 {"t":"Five of ten classes stay below F1 50. The thesis attributes this to few annotations (tires, big bags), high intra-class variance (wood) and inter-class similarity in RGB.","sz":11.5},
 {"t":"The classes that hold are shape- or texture-distinctive (rubble, bulky items, containers); the ones that fail are those where RGB appearance is ambiguous.","sz":11.5},
 {"t":"Per-category behaviour, not the aggregate score, is where the margin is. This is the starting point of the proposal.","b":True,"sz":11.5},
],top=1.15,left=0.45,w=4.10,h=3.90)
img(s,os.path.join(FIG,"alari_f1.png"),4.70,1.15,4.90,3.90)
foot(s,"Numbers: Alari 2024, Table 4.13 (IDA + ResNet50, ten categories), politesi 10589/230633")

# 19 ── synthesis
s=slide("State of the art"); title(s,"The in-scope state of the art at a glance")
table(s,[
 ["Work","Input","Task","Material info"],
 ["Fraternali 2024 (survey)","-","survey, 50 works","declares the material gap"],
 ["AerialWaste 2023","aerial RGB","dataset","22 tags, metadata only"],
 ["Gibellini 2025","aerial RGB","classification","none (presence)"],
 ["Alari 2024","satellite RGB","multi-label cls.","13 categories, wF1 59-69"],
 ["AW ensemble 2025 *","aerial RGB","classification","none (binary)"],
 ["Sun 2023 / CascadeDumpNet 2024","VHR RGB","object detection","none (sites)"],
 ["Disaitek 2024 (vendor)","Pléiades Neo","detection","coarse type qualification"],
 ["CWLD 2024","GF-2 + GE","segmentation","one material (C&D)"],
 ["Saba / Bonifazi / Abbasi / Cilia","WV-3, aerial, HSI","pixel cls.","one material (asbestos)"],
 ["Tanks / vehicles / scrap works","VHR, close range","object detection","shape, not composition"],
],left=0.45,top=1.08,w=9.10,colw=[2.90,1.75,1.95,2.50],fs=8.5,rh=0.38,hh=0.34)
foot(s,"Ten lines summarise every in-scope line of work. Material-level evidence: one multi-label predecessor + one material studied in isolation.")

# 20 ── rgb limits + vnir
s=slide("State of the art"); title(s,"Where RGB falls short, and what VNIR adds")
body(s,[
 {"t":"Different materials share the same colour at 0.3-1.3 m: plastic sheets, asbestos-cement and concrete all appear grey. The separation lies beyond 700 nm (spectral libraries: Kokaly 2017, Knaeps 2020).","sz":11.5},
 {"t":"In the group predecessor, moving from 5 to 10 material categories costs 9.8 F1 points (69.2 to 59.4, Alari 2024). Finer material distinctions are the hard part.","sz":11.5},
 {"t":"Red Edge and NIR separate vegetation, bare soil and weathered surfaces; in Saba 2026 they drive asbestos discrimination.","sz":11.5},
 {"t":"On 10 construction-and-demolition materials, RGB gives 0.87 accuracy; RGB plus two bands at 650-750 and 850-1000 nm gives 0.96, matching the full spectrum (Vitek 2025, laboratory study).","sz":11.5},
 {"t":"Whether this helps waste materials at task GSD, and which ones, has not been measured. It is the object of this thesis.","sz":11.5},
],top=1.10,left=0.45,w=3.85,h=4.0)
img(s,os.path.join(FIG,"vnir_signatures.png"),4.45,1.20,5.15,3.85)
foot(s,"Reflectance: USGS splib07a (Kokaly et al. 2017), 400-1050 nm. Band centres: Maxar, Airbus.")

# 21 ── gaps
s=slide("State of the art"); title(s,"What is missing in the literature")
body(s,[
 {"t":"1.  Multi-material classification has one direct precedent, with ample margin: wF1 59-69 (Alari 2024)."},
 {"t":"2.  No work measures the added value of VNIR bands over RGB for waste materials at very high resolution."},
 {"t":"3.  Results are reported as aggregate scores; per-material behaviour is not analysed."},
 {"t":"4.  Generalisation across regions is rarely evaluated. At site level it costs 5 F1 points (Gibellini 2025)."},
 {"t":"5.  Asbestos is studied on roofs, in isolation, and never inside a waste-material taxonomy."},
])
foot(s,"Each gap maps to one element of the proposal in the next section.")

# ── Part IV: proposed work ────────────────────────────────────────────────

# 22 ── proposal approach
s=slide("Proposed work"); title(s,"Proposed work: approach")
body(s,[
 {"t":"Multi-label image classification, continuing the group line (Gibellini 2025, Alari 2024). Same taxonomy, input extended from RGB to VNIR.","b":True,"sz":12},
 {"t":"Band ablation with everything else fixed: same architecture, same splits, only the input changes.","sz":12},
],top=0.95,h=1.15,mid=False)
img(s,os.path.join(FIG,"pipeline.png"),0.55,2.15,8.9,2.75)
foot(s,"Backbone: Swin-T with remote-sensing pretraining (group baseline); extra bands enter via the input layer.")

# 23 ── imagery
s=slide("Proposed work"); title(s,"Available imagery")
img(s,os.path.join(FIG,"bands_chart.png"),0.55,1.05,8.9,2.85)
body(s,[
 {"t":"Panchromatic: 0.31 m (WorldView-3), 0.30 m (Pléiades Neo). Access: commercial, or free quota via ESA proposal.","sz":12},
 {"t":"This is the band budget of the study: no SWIR, no Sentinel-2. RGB baseline runs on AerialWaste (20-50 cm).","sz":12},
],top=4.05,h=1.05,mid=False)
foot(s,"Band ranges: Maxar WorldView-3, Airbus Pléiades Neo data sheets")

# 24 ── pilot
s=slide("Proposed work"); title(s,"First step: the asbestos pilot")
body(s,[
 {"t":"Asbestos is the entry material, for the reasons seen in the analysis (slide 15): a single material, public pixel-accurate labels, direct regulatory value.","sz":12.5},
 {"t":"1.  Extract roof polygons from the Lombardy WFS registry (10,903 mapped asbestos-cement roofs) and pair them with the available imagery.","sz":12.5},
 {"t":"2.  Build matched RGB and VNIR inputs for the same roofs; add negative roofs from regional building footprints.","sz":12.5},
 {"t":"3.  Train the same classifier on both inputs; compare F1, precision, recall on held-out areas.","sz":12.5},
 {"t":"4.  Decision point: the measured VNIR delta on a single, well-labelled material tells whether the full multi-label extension is worth the acquisition cost.","b":True,"sz":12.5},
])
foot(s,"Ground truth: Regione Lombardia WFS, Mappatura 2020 (EPSG:32632)")

# 25 ── evaluation
s=slide("Proposed work"); title(s,"Proposed work: evaluation")
body(s,[
 {"t":"Per-material F1 alongside weighted and macro averages: aggregates hide exactly the classes this thesis targets.","sz":12},
 {"t":"Delta versus the RGB baseline for each band configuration, with confidence intervals over repeated runs.","sz":12},
 {"t":"Generalisation: train and test on disjoint geographic areas, besides the standard split.","sz":12},
 {"t":"Reference points: wF1 69.2 / 59.4 (Alari 2024); Macro-F1 97.6 (Saba 2026, per-pixel) as upper reference for the pilot, not directly comparable.","sz":12},
 {"t":"If VNIR does not help a material, that is a documented negative result with practical value for sensor choice.","sz":12},
],top=1.15,left=0.45,w=4.7,h=3.9)
img(s,os.path.join(FIG,"eval_grid.png"),5.35,1.55,4.25,2.9)
foot(s,"Metrics: per-class F1, macro-F1, weighted F1, confusion matrices. Splits frozen before testing.")

# 26 ── references 1: the waste-detection line
s=slide(); title(s,"References (1/3): waste detection")
body(s,[
 {"t":"Alari 2024. Fighting environmental crime with deep learning: classifying waste materials from illegal landfills in satellite imagery. M.Sc. thesis, PoliMi, 10589/230633.","sz":10.5},
 {"t":"Fraternali et al. 2024. Solid waste detection, monitoring and mapping in remote sensing images: a survey. Waste Management / arXiv:2402.09066.","sz":10.5},
 {"t":"Gibellini et al. 2025. A deep learning pipeline for solid waste detection in remote sensing images. Waste Management Bulletin.","sz":10.5},
 {"t":"Torres, Fraternali 2023. AerialWaste: a dataset for illegal landfill discovery in aerial images. Scientific Data 10:63.","sz":10.5},
 {"t":"Sharmily et al. 2025. AerialWaste ensembles (preprint).  ·  Zhang, Ma 2024. CascadeDumpNet. Remote Sensing of Environment 313.","sz":10.5},
 {"t":"Sun et al. 2023. Revealing influencing factors on global waste distribution. Nature Communications 14:1444.","sz":10.5},
 {"t":"CWLD 2024. Construction waste landfill dataset. Scientific Data.  ·  Disaitek 2024, vendor report (Airbus).","sz":10.5},
 {"t":"Tisza et al. 2023. Waste detection and change analysis from multispectral satellite imagery (Sentinel-2).  ·  MARIDA 2022. Marine debris benchmark. PLOS ONE.","sz":10.5},
])
foot(s,"Full annotated library: 47 papers, papers/INDEX.md")

# 27 ── references 2: materials and spectra
s=slide(); title(s,"References (2/3): materials and spectra")
body(s,[
 {"t":"Saba et al. 2026. Asbestos-cement roofs from WorldView-3 VNIR. J. Hazardous Materials.  ·  Bonifazi et al. 2026. Geomatics 6(3):41.","sz":10},
 {"t":"Abbasi et al. 2024. Multi-temporal asbestos change detection. RSASE.  ·  Cilia et al. 2015. ISPRS IJGI 4(2).","sz":10},
 {"t":"Shepherd et al. 2025. Asbestos-cement rooftops with EnMAP hyperspectral data. Scientific Reports.","sz":10},
 {"t":"Asbestos slate from drone imagery, 2023 (training-data study).","sz":10},
 {"t":"Vitek et al. 2025. Critical wavelengths for construction and demolition waste materials. Resources, Conservation and Recycling.","sz":10},
 {"t":"Kokaly et al. 2017. USGS Spectral Library v7 (splib07a). USGS DS 1035.  ·  Knaeps et al. 2020. Hyperspectral plastics library. ESSD.","sz":10},
 {"t":"EMIT 2025. Global-scale plastic detection from space. Geophys. Res. Lett.  ·  Zhou et al. 2021. Plastic classifier. Remote Sensing of Environment.","sz":10},
 {"t":"Guo, Li 2020. Normalized difference plastic index on WorldView-3. ISPRS JPRS.  ·  Aguilar et al. 2025. Urban macroplastics with WorldView-3.  ·  Aguilar et al. 2021. WV-3 VNIR/SWIR ablation. Remote Sensing.","sz":10},
 {"t":"SpectralWaste 2024. Multimodal spectral sorting dataset.  ·  Plastic UAV-SWIR segmentation, 2026. Remote Sensing.  ·  Plastic UAV-IoT detection, 2025. J. Hazardous Materials Advances.","sz":10},
])
foot(s,"Spectral-library and laboratory studies cited with their scope stated on the slides that use them")

# 28 ── references 3: objects, platforms, excluded backbones
s=slide(); title(s,"References (3/3): objects, platforms, excluded backbones")
body(s,[
 {"t":"Ramachandran et al. 2024. Deep learning to map well pads and storage tanks. Nature Communications.","sz":10},
 {"t":"YOLOv7-OT 2024. Storage tank detection. Remote Sensing.  ·  Truck-and-container detection from satellite imagery, 2025 (preprint).","sz":10},
 {"t":"Vehicle detection with domain adaptation in VHR imagery, 2020.  ·  ELV Hybrid-YOLOv5 2025. Non-ferrous metals in end-of-life vehicles.","sz":10},
 {"t":"UAV solid-waste segmentation, 2024.  ·  C&D debris volume from UAV photogrammetry, 2022. Drones.  ·  AI-powered drones in waste management: systematic review, 2025. Drones.","sz":10},
 {"t":"fCLIPSeg 2025. Event-agnostic debris segmentation.","sz":10},
 {"t":"DOFA 2024 (Xiong et al.).  ·  AnySat 2025.  ·  Prithvi-EO-2.0 2024 (Szwarcman et al.).  ·  SpectralGPT 2024. TPAMI.","sz":10},
 {"t":"SatMAE 2022 (Cong et al.). NeurIPS.  ·  SSL4EO-S12 2023 (Wang et al.).  ·  SoftCon 2024 (Wang et al.).","sz":10},
 {"t":"Corley et al. 2024. Revisiting pre-trained remote sensing benchmarks.  ·  Foundation models for remote sensing: a survey, 2024.  ·  DEFLECT 2025 (Thoreau et al.). ICCV.","sz":10},
])
foot(s,"47 works cited in total; foundation-model line reported only as an explicit exclusion (slide 10)")

# page numbers
for pos,sl in enumerate(prs.slides,1):
    if pos==1: continue
    tb=sl.shapes.add_textbox(IN(0.10),IN(5.26),IN(0.45),IN(0.28))
    tf=tb.text_frame; tf.word_wrap=False
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=str(pos)
    r.font.name=F; r.font.size=Pt(9); r.font.color.rgb=GREY

prs.save(OUTP)
dl=os.path.expanduser("~/Downloads")
if os.path.isdir(dl):
    import shutil; shutil.copy(OUTP, os.path.join(dl,"slide_v8.pptx"))
print("saved", OUTP, len(prs.slides._sldIdLst), "slides")
