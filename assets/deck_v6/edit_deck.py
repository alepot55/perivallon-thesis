"""Edit ultime_slide.pptx in place: keep the good originals + their embedded
paper images, apply the WV-3 + Pléiades Neo pivot, add slides in the SAME
minimal B/W standard style. Output: assets/deck_v6/deck_bw.pptx (+ ~/Downloads)."""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC=os.path.expanduser("~/Downloads/ultime_slide.pptx")
OUT=os.path.expanduser("/home/alepot55/Desktop/uni/Tesi/assets/deck_v6/deck_bw.pptx")
BW=os.path.expanduser("/home/alepot55/Desktop/uni/Tesi/assets/deck_v6/figs_bw")
CL=os.path.expanduser("/home/alepot55/Desktop/uni/Tesi/assets/deck_v6/figs_color")
INK=RGBColor(0x1A,0x1A,0x1A); GREY=RGBColor(0x66,0x66,0x66)
HDR=RGBColor(0xDD,0xDD,0xDD); ALT=RGBColor(0xF5,0xF5,0xF5); WHITE=RGBColor(0xFF,0xFF,0xFF)
F="Calibri"
prs=Presentation(SRC)
S=list(prs.slides)
EMU=914400
def IN(v): return Inches(v)

# ---- low-level helpers -----------------------------------------------------
def fill_tf(tf, specs, wrap=True):
    """specs: list of dicts {t, b(old), i(talic), sz, c(olor)}."""
    tf.clear(); tf.word_wrap=wrap
    for k,sp in enumerate(specs):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph()
        p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(sp.get("sa",7)); p.line_spacing=sp.get("ls",1.22)
        r=p.add_run(); r.text=sp["t"]
        r.font.name=F; r.font.size=Pt(sp.get("sz",14)); r.font.bold=sp.get("b",False)
        r.font.italic=sp.get("i",False); r.font.color.rgb=sp.get("c",INK)

def boxes(slide):
    """return (title_tb, body_tbs[], footer_tb) by first-run size heuristics."""
    title=foot=None; body=[]
    for sh in slide.shapes:
        if not (getattr(sh,"has_text_frame",False) and sh.has_text_frame): continue
        sz=None; 
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                sz=r.font.size.pt if r.font.size else None; break
            if sz is not None: break
        if sz and sz>=20: title=sh
        elif sz and sz<=9: foot=sh
        else: body.append(sh)
    return title, body, foot

def set_title(slide, text):
    t,_,_=boxes(slide)
    if t: fill_tf(t.text_frame,[{"t":text,"b":True,"sz":23,"c":INK}])

def set_footer(slide, text):
    _,_,f=boxes(slide)
    if f: fill_tf(f.text_frame,[{"t":text,"sz":8,"c":GREY}])

def set_body(slide, specs, idx=0):
    _,b,_=boxes(slide)
    if b and idx<len(b): fill_tf(b[idx].text_frame,specs)

def del_shape(sh): sh._element.getparent().remove(sh._element)

def reencode_all_pictures():
    """PIL-normalize every embedded PNG so LibreOffice renders all slides."""
    import io
    from PIL import Image
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sl in prs.slides:
        for sh in list(sl.shapes):
            if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
                try:
                    im=Image.open(io.BytesIO(sh.image.blob))
                    buf=io.BytesIO(); im.save(buf,"PNG"); buf.seek(0)
                    L,T,W,H=sh.left,sh.top,sh.width,sh.height
                    del_shape(sh); sl.shapes.add_picture(buf,L,T,W,H)
                except Exception as e:
                    print("  reencode skip:",e)

def replace_pic(slide, path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    for sh in list(slide.shapes):
        if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
            L,T,Wb,Hb=sh.left,sh.top,sh.width,sh.height
            del_shape(sh)
            iw,ih=Image.open(path).size; ar=iw/ih
            W=Wb; H=int(W/ar)
            if H>Hb: H=Hb; W=int(H*ar)
            slide.shapes.add_picture(path, L+(Wb-W)//2, T+(Hb-H)//2, width=W, height=H)
            return True
    return False

def replace_biggest_pic(slide, path):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    pics=[sh for sh in slide.shapes if sh.shape_type==MSO_SHAPE_TYPE.PICTURE]
    if not pics: return False
    sh=max(pics, key=lambda x: x.width*x.height)
    L,T,Wb,Hb=sh.left,sh.top,sh.width,sh.height
    del_shape(sh)
    iw,ih=Image.open(path).size; ar=iw/ih
    W=Wb; H=int(W/ar)
    if H>Hb: H=Hb; W=int(H*ar)
    slide.shapes.add_picture(path, L+(Wb-W)//2, T+(Hb-H)//2, width=W, height=H)
    return True

def add_slide():
    s=prs.slides.add_slide(prs.slide_layouts[10])  # BLANK
    # strip any leftover placeholders
    for ph in list(s.placeholders):
        try: del_shape(ph)
        except: pass
    return s

def n_title(s,text):
    tb=s.shapes.add_textbox(IN(0.45),IN(0.34),IN(9.10),IN(0.68))
    fill_tf(tb.text_frame,[{"t":text,"b":True,"sz":23,"c":INK}]); return tb
def n_body(s,specs,top=1.35,w=9.05,h=3.65,vc=True):
    tb=s.shapes.add_textbox(IN(0.5),IN(top),IN(w),IN(h))
    fill_tf(tb.text_frame,specs)
    if vc: tb.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
    return tb
def n_foot(s,text):
    tb=s.shapes.add_textbox(IN(0.45),IN(5.29),IN(9.10),IN(0.26))
    fill_tf(tb.text_frame,[{"t":text,"sz":8,"c":GREY}]); return tb
def n_img(s,path,left=5.30,top=1.25,w=4.40,h=3.85):
    from PIL import Image
    iw,ih=Image.open(path).size; ar=iw/ih
    W=IN(w); H=int(W/ar)
    if H>IN(h): H=IN(h); W=int(H*ar)
    s.shapes.add_picture(path,IN(left)+ (IN(w)-W)//2, IN(top)+(IN(h)-H)//2, width=W,height=H)

def style_table(tbl, headerbold=True, fs=10.5):
    nr=len(tbl.rows); nc=len(tbl.columns)
    for i in range(nr):
        for j in range(nc):
            c=tbl.cell(i,j); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=IN(0.07); c.margin_right=IN(0.05); c.margin_top=IN(0.02); c.margin_bottom=IN(0.02)
            c.fill.solid(); c.fill.fore_color.rgb = HDR if i==0 else (ALT if i%2==0 else WHITE)
            for p in c.text_frame.paragraphs:
                p.line_spacing=1.0
                for r in p.runs:
                    r.font.name=F; r.font.size=Pt(fs); r.font.color.rgb=INK
                    r.font.bold=(i==0 and headerbold)

def new_table(s, data, left=0.5, top=1.35, w=9.05, colw=None, fs=10, rh=0.5, hh=0.42, vcenter=False):
    nr=len(data); nc=len(data[0])
    th=hh+(nr-1)*rh
    if vcenter: top=1.35+(3.65-th)/2.0
    gt=s.shapes.add_table(nr,nc,IN(left),IN(top),IN(w),IN(th)).table
    if colw:
        for j,cw in enumerate(colw): gt.columns[j].width=IN(cw)
    gt.rows[0].height=IN(hh)
    for i in range(1,nr): gt.rows[i].height=IN(rh)
    for i,row in enumerate(data):
        for j,val in enumerate(row): gt.cell(i,j).text=str(val)
    style_table(gt,fs=fs); return gt

# ════════ EDIT EXISTING SLIDES ════════
# Reframe / update kept original slides
set_body(S[1],[
 {"t":"Illegal waste dumping (discariche abusive) is an environmental crime and a public-health issue.","sz":14},
 {"t":"Agencies such as ARPA cannot inspect everything: with limited inspectors, they must prioritise - and priority depends on WHAT is dumped, the material.","sz":14},
 {"t":"Automated pipelines now find sites well (detection is essentially mature) - but they classify presence, not what is there.","sz":14},
 {"t":"The advisor's own survey (Fraternali et al. 2024, 50 works, 1987-2023): nearly all RGB; material identification declared an open gap requiring <=30 cm imagery + multispectral data.","sz":14},
 {"t":"Research question","b":True,"sz":14},
 {"t":"What is the added value of multispectral data, vs. RGB only, for waste material classification from satellite imagery?","sz":14},
])
set_footer(S[1],"Framing distilled from Fraternali et al. 2024 (PoliMi survey, arXiv:2402.09066)")

set_body(S[2],[
 {"t":"VHR aerial / satellite tile (RGB)  →  CNN or Transformer backbone  →  pretraining (ImageNet, EO SSL)  →  two-step fine-tune  →  binary waste / no-waste.","sz":13},
 {"t":"•  Accurate within the geographic context it was trained on: Gibellini 2025, F1 92.0% at 20 cm.","sz":13},
 {"t":"•  Classifies presence, not material.","sz":13},
 {"t":"•  Tied to chromatic information only.","sz":13},
 {"t":"•  Generalizzazione collapses outside it: -5.1% F1 cross-region; worse cross-sensor.","sz":13},
 {"t":"These limits motivate looking at the spectrum.","i":True,"sz":12,"c":GREY},
])
set_footer(S[2],"Gibellini et al. 2025 (Waste Mgmt Bull.) · Fraternali et al. 2024 (survey)")
set_title(S[6],"RGB fails in two distinct ways")
set_body(S[15],[
 {"t":"2024-2025: a wave of pretrained backbones for EO (Prithvi, SpectralGPT, DOFA, AnySat).","sz":14},
 {"t":"Most are pretrained at 10-30 m on Sentinel-2 / HLS - naive transfer to sub-metre VHR is not guaranteed.","sz":14},
 {"t":"Two routes forward","b":True,"sz":13},
 {"t":"•  Sensor-agnostic models (DOFA, AnySat), accept arbitrary band sets.","sz":13},
 {"t":"•  Parameter-efficient adapters (DEFLECT), freeze the backbone, add small heads (<1% parameters).","sz":13},
 {"t":"Open question: does pretraining at 10-30 m transfer to ~1 m VHR - and to material classification?","i":True,"sz":12,"c":GREY},
])

set_body(S[3],[
 {"t":"Each pixel is a vector of reflectance values, the material's spectral signature.","sz":14},
 {"t":"•  RGB,  3 broad visible bands (colour only).","sz":14},
 {"t":"•  Multispectral (MS),  4-15 chosen bands (e.g. WorldView-3 8 VNIR + 8 SWIR; Pleiades Neo 6 VNIR).","sz":14},
 {"t":"•  Hyperspectral (HSI),  hundreds of contiguous narrow bands (EnMAP 230; USGS splib07a 2,151 ch).","sz":14},
 {"t":"This vector is the raw input from which any classifier reasons.","sz":14},
])
# S12 → reframe as the VHR / SWIR landscape (NOT "my data")
set_title(S[11],"Very-high resolution: the SWIR divide")
for sh in list(S[11].shapes):
    if getattr(sh,"has_table",False): del_shape(sh)
new_table(S[11],[
 ["","WorldView-3","Pleiades Neo","SuperDove"],
 ["VNIR","8 b @ 1.24 m","6 b @ 1.2 m","8 b @ 3 m"],
 ["SWIR","8 b @ 3.7 m","none","none"],
 ["Pan","~0.31 m","~0.30 m","-"],
 ["Access","commercial / TPM","commercial / TPM","free, near-daily"],
], left=0.5, top=1.35, w=8.5, colw=[1.4,2.4,2.4,2.3], fs=11, rh=0.5)
_,_b11,_=boxes(S[11])
for _bb in _b11: _bb.top=IN(3.95); _bb.height=IN(1.25)
set_body(S[11],[
 {"t":"WorldView-3 is the only sub-2 m platform that also carries SWIR - where the chemistry lives.","sz":12.5},
 {"t":"Pleiades Neo and SuperDove are VNIR-only: sharp, but blind to the diagnostic SWIR absorptions.","sz":12.5},
 {"t":"Precision (the Sentinel-2 trap): S-2 DOES carry SWIR, at 20 m - its limit here is resolution, not the band. The VHR question is who carries SWIR at all.","sz":12.5,"i":True},
])
set_footer(S[11],"Specs: Maxar WorldView-3 · Airbus Pleiades Neo · Planet SuperDove")
# S13 (was Aguilar 2021) → REPURPOSED: evidence-timeline (recency) slide
set_title(S[12],"A moving field: material-level evidence is 2024-2026")
set_body(S[12],[
 {"t":"Site-level RGB detection matured in 2023-2024; material-level discrimination is being demonstrated NOW (2024-2026) - but not yet inside an illegal-waste pipeline at task resolution.","sz":13},
])
replace_pic(S[12], os.path.join(CL,"evidence_timeline.png"))
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST12
for _sh in list(S[12].shapes):
    if _sh.shape_type==_MST12.PICTURE:
        _sh.top=IN(2.0); _sh.left=IN(1.2); _sh.width=IN(7.6); _sh.height=IN(3.15)
set_footer(S[12],"Years/GSD from the surveyed works · shaded band = coarser than the task (>5 m GSD)")
# S14 more-bands
set_body(S[13],[
 {"t":"•  Vitek 2025 (lab HSI, C&D waste): RGB 0.87 → RGB + 2 NIR bands 0.96 - on par with the full 768-band HSI.","sz":14},
 {"t":"•  SpectralWaste 2024: RGB+SWIR fusion beats either alone (mIoU 58 vs 48); SWIR wins the thin, iso-colour classes.","sz":14},
 {"t":"•  Historical anchor - Aguilar 2021 (WV-3): VNIR 90.85 → All 97.38, and the jump came from SWIR.","sz":14},
])
_tb13=S[13].shapes.add_textbox(IN(0.45),IN(4.82),IN(9.10),IN(0.45))
fill_tf(_tb13.text_frame,[{"t":"Why it matters: WHICH bands is material-dependent - VNIR suffices on some materials (Saba 2026); chemistry-bound hazards need SWIR. Exactly what a controlled ablation must measure.","i":True,"sz":11,"c":GREY}])
# S17 DOFA
set_body(S[16],[
 {"t":"Dynamic One-For-All (Xiong et al. 2024).","b":True,"sz":14},
 {"t":"•  A hypernetwork generates patch-embedding weights from each band's central wavelength.","sz":14},
 {"t":"•  One backbone ingests any sensor (WV-3 16 b, Pleiades Neo 6 b, Sentinel-2 13 b, HSI) by wavelength.","sz":14},
 {"t":"•  Makes an RGB vs VNIR vs SWIR comparison fair - the substrate for a controlled band ablation.","sz":14},
 {"t":"•  Pairable with PEFT adapters (DEFLECT, <1% params) for cheap spectral fine-tuning.","sz":14},
])
# S18 gaps (expanded to the SOTA-revealed gaps)
set_title(S[17],"Gaps the state of the art reveals")
set_body(S[17],[
 {"t":"•  RGB ceiling: every operational VHR waste detector is RGB-only and learns shape/context, not material.","sz":13},
 {"t":"•  No public dataset combines VHR + terrestrial waste + material labels (AerialWaste is RGB; CWLD is C&D-only).","sz":13},
 {"t":"•  Material evidence is adjacent, not direct: proven on greenhouses / asbestos roofs in isolation, never in a waste pipeline.","sz":13},
 {"t":"•  Generalization (cross-region / cross-sensor / cross-time) is fragile and unmeasured for MS.","sz":13},
 {"t":"•  Foundation models pretrain at 10-30 m and are untested for material at VHR.","sz":13},
 {"t":"•  No link from material → hazard class → risk, the decision output agencies actually need.","sz":13},
])
set_footer(S[17],"Gaps distilled across the surveyed works · Fraternali et al. 2024")

replace_pic(S[8], os.path.join(CL,"sensor_radar.png"))
replace_pic(S[13], os.path.join(CL,"bands_plateau.png"))
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST13
for _sh in list(S[13].shapes):
    if _sh.shape_type==_MST13.PICTURE:
        _sh.top=IN(2.68); _sh.left=IN(2.10); _sh.width=IN(5.8); _sh.height=IN(2.05)
replace_biggest_pic(S[6], os.path.join(CL,"rgb_fails.png"))
# S15 SuperDove → Honest caveat (+ swir8 figure)
sd=S[14]
set_title(sd,"Honest limit: resolution is split between texture and chemistry")
_,bd,_=boxes(sd)
if bd:
    bd[0].left=IN(0.45); bd[0].top=IN(1.25); bd[0].width=IN(4.55); bd[0].height=IN(3.9)
    fill_tf(bd[0].text_frame,[
     {"t":"•  Spectra are sampled coarsely: VNIR ~1.24 m, SWIR ~3.7 m (~14 m2 pixels).","sz":13},
     {"t":"•  The pan band gives fine texture ~0.3 m, but no material info alone.","sz":13},
     {"t":"•  A small dump is easier to SEE (texture) than to chemically IDENTIFY (spectra).","sz":13},
     {"t":"•  SWIR-8 bottleneck: asbestos 2.32 + concrete 2.34 + plastic 2.31 crowd one band.","sz":13},
     {"t":"•  And SWIR is not always needed: Saba 2026 / Abbasi 2024 classify asbestos at VHR with none.","sz":13},
    ])
n_img(sd, os.path.join(CL,"swir8_bottleneck.png"), left=5.05, top=1.3, w=4.7, h=3.7)
set_footer(sd,"USGS splib07a (Kokaly et al., 2017) · WV-3 SWIR band centres (Maxar)")

# ════════ NEW SLIDES ════════
TCW=[2.15,2.25,2.35,2.35]  # survey-table column widths (sum 9.10)
def survey(title, rows, foot):
    n=add_slide(); n_title(n,title)
    new_table(n,[["Work (year)","Input / GSD","Task","Method","Key result"]]+rows,
              colw=[1.72,1.62,1.28,1.98,2.45], fs=9, rh=0.52, vcenter=True)
    n_foot(n,foot); return n

# N: risk (motivation)
n=add_slide(); n_title(n,"Why material matters: risk = hazard x exposure x magnitude")
n_body(n,[
 {"t":"•  The hazard is material-bound: inert rubble vs plastics vs asbestos-cement are worlds apart.","sz":14},
 {"t":"•  Material carries the hazard: the European Waste Catalogue (EWC) flags hazardous codes with * - asbestos-containing construction material is 17 06 05*.","sz":14},
 {"t":"•  Priority risk = hazard (material) x exposure (who/what is nearby) x magnitude (how much).","sz":14},
 {"t":"•  Regione Lombardia already operationalises this for asbestos roofs: the Indice di Degrado (d.d.g. 13237/2008) score triggers re-evaluation, 3-year remediation, or 12-month removal.","sz":14},
 {"t":"•  So identifying the MATERIAL - not just locating the site - is the decision-relevant question.","b":True,"sz":14},
])
n_foot(n,"EWC List of Waste (Dec. 2000/532/EC) · Indice di Degrado d.d.g. 13237/2008 · Fazzo et al. 2023")

# G1 RGB detection
survey("State of the art: RGB waste detection",
 [["Gibellini 2025","AerialWaste RGB 20 cm","classification","Swin-T + RSP, two-step","F1 92.0%; cross-region -5.1%"],
  ["Disaitek 2024","Pleiades Neo 0.3 m","detection","operational service","waste >=2 m2 ~95% (vendor)"],
  ["CascadeDumpNet 2024","Pleiades 0.5 m RGB","object detection","cascade CNN + AutoML","mAP 84.6%; transfers across cities"],
  ["Sun 2023 (global)","VHR RGB 0.3-1 m","object detection","BCA-Net (Faster R-CNN)","~2,500 dumpsites; sens. 98%"],
  ["AerialWaste (Torres 2023)","aerial RGB 20-50 cm","dataset + cls.","ResNet+FPN baseline","F1 80.7%; 22 categories, RGB-only"]],
 "All RGB: detect site shape/context; material composition stays out of reach. † vendor/pre-print not peer-benchmarked.")

# band->material (physics)
n=add_slide(); n_title(n,"Where each hazard becomes separable: feature → band")
n_img(n, os.path.join(CL,"band_material_map.png"), left=0.6, top=1.15, w=8.8, h=3.95)
n_foot(n,"Diagnostic absorptions RGB cannot see · USGS splib07a (Kokaly 2017) · Aguilar 2021/2025 · Cilia 2015")

# G2 asbestos
survey("State of the art: asbestos discrimination",
 [["Saba 2026 †","WV-3 VNIR 1.24 m","pixel classification","32 classifiers, Fine-KNN","Macro-F1 97.6% (VNIR only)"],
  ["Bonifazi 2026 †","WV-3 16 b 1.24/3.7 m","pixel cls. + aggregation","Max-Likelihood, multi-temporal","AC roofs; removal monitoring"],
  ["Shepherd 2025","EnMAP HSI 30 m","target detection","8-classifier cascade (ACE)","86% field match; OA 91.4%"],
  ["Abbasi 2024 †","aerial RGB (no SWIR)","OBIA classification","DenseNet + LSTM, temporal","OA ~96%, AC 94% by shape+time"],
  ["Cilia 2015 (hist.)","MIVIS airborne HSI 3 m","pixel classification","SAM + MNF, ISD index","PA 89% / UA 86%"]],
 "Asbestos = the sharpest proof MS/HSI recovers material identity. † paywalled / pre-print - motivating, not established.")

# G3 plastics & urban materials
survey("State of the art: plastics & urban materials",
 [["Aguilar 2025","WV-3 SWIR 3.7 m","target detection","matched filter (USGS spectra)","precision 92.5%; lab-image r=0.95"],
  ["CDW (Vitek) 2025","lab HSI 768 b","classification","narrowband selection + MLP","RGB 0.87 → +2 NIR 0.96 ~ full HSI"],
  ["EMIT (Estrela 2025)","EMIT HSI 60 m","target detection","matched filter, HDPE/PVC","first global orbital plastic map"],
  ["SpectralWaste 2024","RGB + SWIR HSI (conveyor)","segmentation","RGB/HSI/fusion (CMX)","fusion mIoU 58 > RGB 48"],
  ["MARIDA 2022","Sentinel-2 10 m","segmentation","RF / U-Net benchmark","F1 0.79; sub-pixel mixing at 10 m"],
  ["Aguilar 2021 (hist.)","WV-3 VNIR + SWIR","OBIA classification","DT, 3-way band ablation","OA 90.85 → 96.79 → 97.38"]],
 "Material identity comes from NIR-SWIR chemistry (C-H, Mg-OH), not colour - but proven outside illegal-waste imagery.")

# G5 datasets
survey("State of the art: datasets & the data gap",
 [["CWLD 2024","GF-2 0.8 m + GE 0.5 m","segmentation","C&D masks, DeepLabV3+","F1 88.9% / IoU 82% (China)"],
  ["SpectralWaste 2024","RGB + SWIR HSI","segmentation","conveyor-belt, 224-band","sorting plant, not Earth observation"],
  ["AerialWaste (Torres 2023)","aerial RGB 20-50 cm","classification","22 material categories","RGB-only; coordinates withheld"],
  ["MARIDA 2022","Sentinel-2 10 m","segmentation","15-class pixel benchmark","marine debris - water, not land"]],
 "No public dataset combines very-high resolution + terrestrial waste + material labels - the structural gap.")

# G6 foundation models
survey("State of the art: foundation models",
 [["AnySat (Astruc 2025)","11 sensors 0.2-250 m","pretraining","JEPA multimodal","SOTA on 9 tasks"],
  ["DEFLECT 2025","MS 6-13 b","adaptation","PEFT adapter <1% params","on par with full fine-tune"],
  ["DOFA (Xiong 2024)","5 modalities 1-30 m","pretraining","wavelength hypernetwork","202 bands; new-sensor plug-in"],
  ["SoftCon 2024","Sentinel-1/2 13 b","pretraining","soft contrastive + DINOv2","SOTA 10/11; random-init extra bands"],
  ["Prithvi-EO-2.0 2024","HLS 6 b incl. SWIR 30 m","pretraining","3D temporal MAE","+8 pp GEO-Bench; fixed bands"]],
 "Most pretrain at 10-30 m on Sentinel-2/HLS; transfer to VHR and to material classification is untested.")

# G7 object vs material
survey("State of the art: object vs material",
 [["ELV Hybrid-YOLOv5 2025","close-range infrared","object detection","Hybrid-YOLOv5","mAP 84%; needs spectrum / close-range"],
  ["Ramachandran 2024","VHR sub-metre","object detection","DL detector (Nat. Comm.)","tanks P 0.96 / R 0.97; >169k mapped"],
  ["YOLOv7-OT 2024","VHR satellite","object detection","YOLOv7 + CBAM","tanks 90% acc / 95.9% precision"],
  ["UAV solid-waste 2024","UAV","segmentation","dual-branch network","OA >94%; generic pile, no per-material"]],
 "Objects (tanks, vehicles) are mature on shape; material composition (scrap, slag) still needs the spectrum.")

# generalization (table)
n=add_slide(); n_title(n,"Generalization: fragile and largely unmeasured for MS")
new_table(n,[
 ["","In-domain","Cross-region","Cross-sensor"],
 ["RGB (today)","strong","-5.1% F1 (Gibellini)","ports trivially"],
 ["+ VNIR","-","gap narrows (physics)","band-shift fragile"],
 ["+ SWIR","-","gap narrows most","widens if not harmonized"],
], left=0.9, top=1.55, w=7.5, colw=[1.8,1.9,1.95,1.85], fs=12, rh=0.6)
n_body(n,[{"t":"A recurring SOTA blind spot: most works report in-domain only. Whether added bands narrow or widen the gap is itself an open question.","i":True,"sz":12,"c":GREY}], top=3.7, h=0.8, vc=False)
n_foot(n,"Cross-region -5.1%: Gibellini 2025 · band-shift fragility: GeoCrossBench, DOFA/AnySat · pattern, not yet measured for this task")

# proposed direction (light, dataset-independent)
n=add_slide(); n_title(n,"What the state of the art points to")
n_body(n,[
 {"t":"The open question is sharp and unanswered: how much does multiband beat RGB for waste MATERIAL - and where does it stop?","b":True,"sz":14},
 {"t":"•  A controlled band ablation (RGB → +RedEdge/NIR → full VNIR → +SWIR) on a wavelength-agnostic backbone (DOFA).","sz":13},
 {"t":"•  Generalization as a first-class axis (cross-region / sensor / time), not an afterthought.","sz":13},
 {"t":"•  Candidate VHR platforms span the SWIR divide: WorldView-3 (with SWIR) vs Pleiades Neo (VNIR-only).","sz":13},
 {"t":"•  Asbestos is the cleanest first material: public ground truth + a textbook SWIR signature.","sz":13},
 {"t":"Independent of any single dataset: the SOTA defines the experiment, not the data in hand.","i":True,"sz":12,"c":GREY},
])
n_foot(n,"Direction grounded in the surveyed gaps · backbone: DOFA (Xiong 2024) · baseline ref.: Gibellini 2025")

# indices of the appended new slides (creation order)
_I_RISK, _I_G1, _I_BM, _I_G2, _I_G3, _I_G5, _I_G6, _I_G7, _I_GEN, _I_DIR = range(len(S), len(S)+10)

# Rebuild the 2 slides LibreOffice drops (fingerprint S[5], sensor S[8])
import io
from PIL import Image as _PILImage
from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
def _extract_pic(slide, outpath):
    for sh in slide.shapes:
        if sh.shape_type==_MST.PICTURE:
            _PILImage.open(io.BytesIO(sh.image.blob)).save(outpath)
            return sh.left, sh.top, sh.width, sh.height
    return None
_fp=os.path.join(BW,"_fingerprint.png"); _sn=os.path.join(BW,"_sensor.png")
_fpg=_extract_pic(S[5], _fp); _sng=_extract_pic(S[8], _sn)
nf=add_slide(); n_title(nf,"Every material has a spectral fingerprint")
n_img(nf, os.path.join(CL,"fingerprint.png"), left=0.55, top=1.10, w=8.9, h=4.0)
n_foot(nf,"USGS splib07a (Kokaly et al., 2017). Diagnostic features: Cilia 2015, Aguilar 2025.")
ns=add_slide(); n_title(ns,"The sensor trade-off: spatial x spectral x revisit")
n_body(ns,[
 {"t":"Three axes, one photon budget. No satellite maximises all three.","sz":12},
 {"t":"•  WorldView-3:  1.24 m VNIR + 3.7 m SWIR, 16 bands - the only VHR + SWIR option.","sz":12},
 {"t":"•  Pleiades Neo:  1.2 m, 6 VNIR, no SWIR - sharpest VNIR-only.","sz":12},
 {"t":"•  Sentinel-2 / EnMAP:  10-30 m - rich spectra, too coarse for material at site scale.","sz":12},
 {"t":"•  SuperDove:  3 m, 8 VNIR, free, near-daily - wide-area but VNIR-only.","sz":12},
], top=1.20, w=4.55)
if _sng: ns.shapes.add_picture(_sn, _sng[0], _sng[1], width=_sng[2], height=_sng[3])
n_foot(ns,"Specs: Planet, ESA, Maxar, ASI, DLR (mission documents).")
_I_FP=len(list(prs.slides))-2; _I_SN=len(list(prs.slides))-1

reencode_all_pictures()

# ════════ REORDER (SOTA narrative) ════════
sldIdLst=prs.slides._sldIdLst
ids=list(sldIdLst)
order=[0, 1, _I_RISK, 2, _I_G1, 3, _I_FP, 6, _I_G2, _I_G3, 12, 13, _I_SN, 11, _I_BM, 14, _I_G5, 15, _I_G6, 16, _I_G7, 17, _I_GEN, _I_DIR]
assert len(order)==24 and all(0<=i<len(ids) for i in order), (len(ids),len(order))
desired=[ids[i] for i in order]
for sid in ids: sldIdLst.remove(sid)
for sid in desired: sldIdLst.append(sid)

# uniform page numbers: bottom-LEFT on every slide except the title
from pptx.enum.text import PP_ALIGN as _AL
for pos,sl in enumerate(prs.slides,1):
    for sh in list(sl.shapes):
        if sh.is_placeholder and sh.placeholder_format.type is not None and "SLIDE_NUMBER" in str(sh.placeholder_format.type):
            del_shape(sh)
    if pos==1: continue
    tb=sl.shapes.add_textbox(IN(0.10),IN(5.28),IN(0.45),IN(0.28))
    tf=tb.text_frame; tf.word_wrap=False
    pr=tf.paragraphs[0]; pr.alignment=_AL.LEFT
    r=pr.add_run(); r.text=str(pos)
    r.font.name=F; r.font.size=Pt(9); r.font.color.rgb=GREY

# normalize ALL bottom footnotes: same box, same size (Calibri 8, grey)
for _sl in prs.slides:
    for _sh in _sl.shapes:
        if not (getattr(_sh,"has_text_frame",False) and _sh.has_text_frame): continue
        _t=_sh.text_frame.text.strip()
        if not _t or _t.isdigit(): continue
        if _sh.top is None or _sh.width is None: continue
        if _sh.top >= IN(5.02) and _sh.width > IN(1.0):
            _sh.left=IN(0.45); _sh.top=IN(5.29); _sh.width=IN(9.10); _sh.height=IN(0.26)
            _sh.text_frame.word_wrap=True
            for _pp in _sh.text_frame.paragraphs:
                _pp.line_spacing=1.0
                for _r in _pp.runs:
                    _r.font.name=F; _r.font.size=Pt(8); _r.font.bold=False
                    _r.font.italic=False; _r.font.color.rgb=GREY

prs.save(OUT)
import shutil; shutil.copy(OUT, os.path.expanduser("~/Downloads/ultime_slide_v2.pptx"))
print("saved", OUT, "·", len(list(prs.slides)), "slides")
