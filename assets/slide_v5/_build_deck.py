#!/usr/bin/env python3
"""
Build Deck v5 — Thomas-style B/N minimal.

20 slides, English, bullet:true (real bullets), no dashes, refs on cited slides.
Embeds assets from assets/slide_v5/ (matplotlib PNGs + cropped PDF pages).

Run: python3 _build_deck.py
Output: papers/decks/deck_v5_thomas.pptx (or local Downloads, configurable)
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path("/home/alepot55/Desktop/uni/Tesi")
ASSETS = ROOT / "assets" / "slide_v5"
PDF_FIGS = ASSETS / "pdf_figs"
OUT_DIR = ROOT / "Downloads_v5"
OUT_DIR.mkdir(exist_ok=True)
OUT = OUT_DIR / "deck_v5_thomas.pptx"

# Style constants
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY_DARK = RGBColor(0x33, 0x33, 0x33)
GREY_MID = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xBB, 0xBB, 0xBB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"  # safer for cross-platform than DejaVu
TITLE_PT = 28
SUBTITLE_PT = 14
BODY_PT = 16
SMALL_PT = 11
REF_PT = 9

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_paragraph(para, text, size=BODY_PT, bold=False, italic=False,
                  color=BLACK, align=PP_ALIGN.LEFT, bullet=False, level=0):
    para.text = ""
    para.alignment = align
    para.level = level
    if bullet:
        # Real round bullet via pPr buChar
        pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
        # Remove any existing bullet
        for tag in ["a:buNone", "a:buChar", "a:buAutoNum"]:
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        bu = etree.SubElement(pPr, qn("a:buChar"))
        bu.set("char", "•")
        # Bullet indent (slight)
        pPr.set("indent", "-228600")  # 0.25 inch hang indent
        pPr.set("marL", "228600")
    else:
        pPr = para._pPr if para._pPr is not None else para._p.get_or_add_pPr()
        for tag in ["a:buChar", "a:buAutoNum"]:
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        # No bullet
        bn = pPr.find(qn("a:buNone"))
        if bn is None:
            etree.SubElement(pPr, qn("a:buNone"))

    run = para.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_title(slide, text, top=Inches(0.35), left=Inches(0.5),
              width=Inches(12.3), height=Inches(0.8)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_paragraph(tf.paragraphs[0], text, size=TITLE_PT, bold=False, color=BLACK)
    return tb


def add_subtitle(slide, text, top=Inches(1.1), left=Inches(0.5),
                 width=Inches(12.3), height=Inches(0.5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_paragraph(tf.paragraphs[0], text, size=SUBTITLE_PT,
                  italic=True, color=GREY_DARK)
    return tb


def add_body(slide, items, top=Inches(1.7), left=Inches(0.5),
             width=Inches(12.3), height=Inches(4.5), size=BODY_PT,
             line_spacing=1.2):
    """items: list of (text, is_bullet, level)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0

    for i, item in enumerate(items):
        if isinstance(item, str):
            text, is_bullet, level = item, True, 0
        else:
            text, is_bullet, level = item
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.line_spacing = line_spacing
        set_paragraph(para, text, size=size, bullet=is_bullet, level=level)
    return tb


def add_refs(slide, refs_text, top=Inches(6.95), left=Inches(0.5),
             width=Inches(12.3)):
    """Small reference line at bottom of slide."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_paragraph(tf.paragraphs[0], refs_text, size=REF_PT,
                  italic=True, color=GREY_MID)
    return tb


def add_page_num(slide, n, total):
    tb = slide.shapes.add_textbox(Inches(12.5), Inches(6.95),
                                  Inches(0.7), Inches(0.4))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_paragraph(tf.paragraphs[0], f"{n} / {total}",
                  size=REF_PT, color=GREY_MID, align=PP_ALIGN.RIGHT)
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        # Add placeholder text instead
        tb = slide.shapes.add_textbox(left, top,
                                       width or Inches(6),
                                       height or Inches(2))
        tf = tb.text_frame
        set_paragraph(tf.paragraphs[0],
                      f"[asset missing: {Path(path).name}]",
                      size=10, italic=True, color=GREY_MID)
        return tb
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, height=height)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


# =========================================================
# Build deck
# =========================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 20

    # --- Slide 1 — Title ---
    s = new_slide(prs)
    # Centered title block
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.2),
                               Inches(11.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0],
                  "Multispectral satellite imagery for",
                  size=32, color=BLACK, align=PP_ALIGN.CENTER)
    p2 = tf.add_paragraph()
    set_paragraph(p2, "illegal waste material classification",
                  size=32, color=BLACK, align=PP_ALIGN.CENTER)
    tb2 = s.shapes.add_textbox(Inches(1.0), Inches(3.9),
                                Inches(11.3), Inches(0.6))
    tf2 = tb2.text_frame
    set_paragraph(tf2.paragraphs[0],
                  "State of the art and thesis direction",
                  size=18, italic=True, color=GREY_DARK,
                  align=PP_ALIGN.CENTER)

    tb3 = s.shapes.add_textbox(Inches(1.0), Inches(5.2),
                                Inches(11.3), Inches(1.4))
    tf3 = tb3.text_frame
    for i, line in enumerate([
        "Alessandro Potenza",
        "M.Sc. Computer Science & Engineering — Artificial Intelligence",
        "Supervisor: Prof. Piero Fraternali",
        "PERIVALLON · Horizon Europe · Grant 101073952",
        "Politecnico di Milano · DEIB · May 2026",
    ]):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        size = 14 if i == 0 else 11
        color = BLACK if i == 0 else GREY_DARK
        bold = i == 0
        set_paragraph(p, line, size=size, color=color, bold=bold,
                      align=PP_ALIGN.CENTER)

    # --- Slide 2 — Problem (riscritta, punti impliciti) ---
    s = new_slide(prs)
    add_title(s, "The problem: priority, not just presence")
    add_body(s, [
        ("Illegal waste dumping is a European environmental crime and a public-health issue.", True, 0),
        ("Agencies such as ARPA cannot inspect everything — they need to prioritise.", True, 0),
        ("Priority depends on the material: asbestos and hazardous waste demand fast intervention, inert rubble does not.", True, 0),
        ("Today's automated pipelines find sites well; they classify presence, not what is there.", True, 0),
    ], top=Inches(1.5), size=BODY_PT)

    # RQ box
    rq_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(5.0),
                                 Inches(12.3), Inches(1.5))
    rq_box.fill.background()
    rq_box.line.color.rgb = BLACK
    rq_box.line.width = Pt(1)
    rq_tf = rq_box.text_frame
    rq_tf.margin_left = rq_tf.margin_right = Inches(0.3)
    rq_tf.margin_top = rq_tf.margin_bottom = Inches(0.15)
    rq_tf.word_wrap = True
    set_paragraph(rq_tf.paragraphs[0], "Research question",
                  size=12, bold=True, color=BLACK)
    pq = rq_tf.add_paragraph()
    set_paragraph(pq,
                  "What is the added value of multispectral data, vs. RGB only, for waste material classification from satellite imagery?",
                  size=15, italic=True, color=BLACK)
    add_page_num(s, 2, TOTAL)

    # --- Slide 3 — Today's paradigm ---
    s = new_slide(prs)
    add_title(s, "Today's paradigm: RGB deep learning on aerial imagery")
    add_image(s, ASSETS / "03_today_paradigm.png",
              Inches(0.7), Inches(1.4), width=Inches(11.9))
    add_body(s, [
        ("Dominant approach reaches strong in-distribution accuracy on RGB.", True, 0),
        ("Limits: classifies presence not material; tied to chromatic information; generalisation drops out-of-area.", True, 0),
        ("These limits motivate looking at the spectrum.", True, 0),
    ], top=Inches(4.5), size=14)
    add_refs(s, "Several recent surveys, e.g. Fraternali et al. 2024.")
    add_page_num(s, 3, TOTAL)

    # --- Slide 4 — What a satellite measures ---
    s = new_slide(prs)
    add_title(s, "A satellite image is a spectral cube, not a photograph")
    add_body(s, [
        ("Each pixel is a vector of reflectance values — the material's spectral signature.", True, 0),
        ("RGB: 3 broad visible bands (colour only).", True, 0),
        ("Multispectral (MS): 4–15 chosen bands. SuperDove records 8.", True, 0),
        ("Hyperspectral (HSI): hundreds of contiguous narrow bands.", True, 0),
        ("This vector is the raw input from which any classifier reasons.", True, 0),
    ], top=Inches(1.5), size=16)
    add_refs(s, "Spectrum schema after USGS splib07a (Kokaly 2017). Band counts from sensor specifications.")
    add_page_num(s, 4, TOTAL)

    # --- Slide 5 (was 4.1) — Three views of one scene ---
    s = new_slide(prs)
    add_title(s, "From bands to information: three views of the same scene")
    add_image(s, ASSETS / "04.1_three_views.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_body(s, [
        ("Same raw spectral cube — many ways to read it.", True, 0),
        ("Each combination highlights a different physical property.", True, 0),
        ("Material discrimination needs more than RGB.", True, 0),
    ], top=Inches(5.6), size=12)
    add_refs(s, "Planet PSScene, Lombardy 2026 (4-band SR product).")
    add_page_num(s, 5, TOTAL)

    # --- Slide 6 (was 5) — Spectral signatures ---
    s = new_slide(prs)
    add_title(s, "Every material has a spectral fingerprint")
    add_image(s, ASSETS / "05_spectral_signatures.png",
              Inches(0.5), Inches(1.3), width=Inches(12.3))
    add_refs(s, "After USGS splib07a (Kokaly 2017) · diagnostic features from Cilia 2015, Zhou 2021.")
    add_page_num(s, 6, TOTAL)

    # --- Slide 7 (was 6) — RGB fails ---
    s = new_slide(prs)
    add_title(s, "RGB fails in two distinct ways")
    add_image(s, ASSETS / "06_rgb_fails.png",
              Inches(0.3), Inches(1.3), width=Inches(12.7))
    add_body(s, [
        ("Iso-chromaticity and sub-pixel mixing — both invisible without spectral information.", True, 0),
        ("Both failure modes vanish in SWIR: features are chemistry of molecular bonds, not albedo.", True, 0),
    ], top=Inches(5.7), size=13)
    add_refs(s, "After USGS splib07a · Zhou 2021 · Aguilar 2025.")
    add_page_num(s, 7, TOTAL)

    # --- Slide 8 (was 7) — Sensor trade-off radar ---
    s = new_slide(prs)
    add_title(s, "The sensor trade-off: spatial × spectral × revisit")
    add_image(s, ASSETS / "07_radar_sensors.png",
              Inches(0.3), Inches(1.2), height=Inches(5.4))
    # Text right side
    tb = s.shapes.add_textbox(Inches(8.0), Inches(1.7),
                               Inches(5.1), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0],
                  "Three axes, one budget. Photon-budget limit, not just price.",
                  size=12, italic=True, color=GREY_DARK)
    p_blank = tf.add_paragraph()
    set_paragraph(p_blank, "", size=6)
    for line in [
        "WorldView-3: best spatial + SWIR; tasked, commercial.",
        "Sentinel-2: free + global; 10 m blur.",
        "PRISMA / EnMAP: hyperspectral; 30 m, slow revisit.",
        "SuperDove: free, 3 m, 8 bands, sub-daily.",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=12, bullet=True, color=BLACK)
    add_refs(s, "Sensor specifications: Planet, ESA, NASA, Maxar, ASI, DLR.")
    add_page_num(s, 8, TOTAL)

    # --- Slide 9 — Survey low-res table ---
    s = new_slide(prs)
    add_title(s, "Multispectral at low resolution (10–30 m)")
    # Table
    rows, cols = 5, 4
    tbl = s.shapes.add_table(rows, cols,
                              Inches(0.5), Inches(1.4),
                              Inches(12.3), Inches(3.5)).table

    # Set column widths
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(2.8)
    tbl.columns[2].width = Inches(3.2)
    tbl.columns[3].width = Inches(2.8)

    headers = ["Work", "Input", "Method", "Result"]
    data = [
        ["MARIDA (Kikaki 2022)", "S-2 13-band @10 m",
         "RF + U-Net, pixel-level", "F1 0.79 multi-class (15)"],
        ["Global dumpsites (Sun 2023)", "VHR commercial",
         "BCA-Net deep learning", "Sens. 98% / Prec. 70%, 28 cities"],
        ["Tisza (Magyar 2023)", "S-2 + PlanetScope",
         "RF + Plastic Index", "Accuracy 96% change det."],
        ["Shepherd 2025 (asbestos)", "EnMAP 230-b @30 m",
         "Cascade 8 classifier", "86% match, 823 detections"],
    ]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        set_paragraph(p, h, size=11, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK

    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            set_paragraph(p, val, size=10, color=BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)

    add_body(s, [
        ("Strong on-binary detection where signal exists in 10–30 m pixels.", True, 0),
        ("Material classification only when hyperspectral SWIR is available (EnMAP).", True, 0),
    ], top=Inches(5.4), size=12)
    add_refs(s, "Kikaki 2022 PLOS ONE · Sun 2023 Nature Comms · Magyar 2023 arXiv · Shepherd 2025 Sci Rep.")
    add_page_num(s, 9, TOTAL)

    # --- Slide 10 — Survey low-res visual ---
    s = new_slide(prs)
    add_title(s, "Low-resolution: what they catch, what they miss")
    add_image(s, PDF_FIGS / "08_marida_map_crop.png",
              Inches(0.4), Inches(1.3), width=Inches(6.3))
    add_image(s, PDF_FIGS / "08_dumpsites_p4_crop.png",
              Inches(7.0), Inches(1.3), height=Inches(4.5))
    # Captions
    tb1 = s.shapes.add_textbox(Inches(0.4), Inches(5.9),
                                Inches(6.3), Inches(0.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    set_paragraph(tf1.paragraphs[0],
                  "MARIDA (Kikaki 2022): 11 coastal sites, Sentinel-2 @10 m. Marine debris confused with natural organic material where pixels mix.",
                  size=10, italic=True, color=GREY_DARK)
    tb2 = s.shapes.add_textbox(Inches(7.0), Inches(5.9),
                                Inches(6.0), Inches(0.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2.paragraphs[0],
                  "Global dumpsites (Sun 2023): BCA-Net on 28 cities. Signal exists at VHR, but material attribution stays out-of-reach.",
                  size=10, italic=True, color=GREY_DARK)
    add_refs(s, "Kikaki et al. 2022 PLOS ONE, Fig. 2 · Sun et al. 2023 Nature Comms 14:1565.")
    add_page_num(s, 10, TOTAL)

    # --- Slide 11 — Survey high-res table ---
    s = new_slide(prs)
    add_title(s, "Multispectral at very-high resolution: WorldView-3")
    rows, cols = 5, 4
    tbl = s.shapes.add_table(rows, cols,
                              Inches(0.5), Inches(1.4),
                              Inches(12.3), Inches(3.5)).table

    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(2.8)
    tbl.columns[2].width = Inches(3.2)
    tbl.columns[3].width = Inches(2.8)

    headers = ["Work", "Input", "Method", "Result"]
    data = [
        ["Aguilar 2021 (ablation)", "WV-3 VNIR+SWIR @1.2 / 3.7 m",
         "OBIA + RF: VNIR / SWIR / All", "OA 90.85 → 96.79 → 97.38 %"],
        ["Aguilar 2025 (macroplastics)", "WV-3 SWIR 8-b @3.7 m",
         "Matched filter + lab spectra", "Precision 92–95 % riparian"],
        ["Zhou 2021 (KB-classifier)", "WV-3 SWIR 8-b",
         "Decision tree knowledge-based", "3 polymer clusters, OA >80%"],
        ["Bonifazi 2026 (asbestos)", "WV-3 VNIR+SWIR @1.2 m",
         "Py6S + MLC + building-level", "F1 0.87 multi-year Mantova"],
    ]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        set_paragraph(p, h, size=11, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            set_paragraph(p, val, size=10, color=BLACK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
    add_body(s, [
        ("VHR + SWIR is today's strongest setup for material classification.", True, 0),
        ("But cost-prohibitive for area-wide coverage (€15–30 / km²).", True, 0),
    ], top=Inches(5.4), size=12)
    add_refs(s, "Aguilar 2021 RS 13:2133 · Aguilar 2025 EMA · Zhou 2021 RSE 264:112598 · Bonifazi 2026 Geomatics.")
    add_page_num(s, 11, TOTAL)

    # --- Slide 12 — Aguilar ablation big chart ---
    s = new_slide(prs)
    add_title(s, "Spectral added value, measured: Aguilar 2021 on WV-3")
    add_image(s, ASSETS / "10_aguilar_ablation.png",
              Inches(0.4), Inches(1.2), height=Inches(4.4))
    # Text on right
    tb = s.shapes.add_textbox(Inches(7.8), Inches(1.4),
                               Inches(5.3), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0],
                  "Canonical benchmark for 'RGB vs multispectral added value'.",
                  size=12, italic=True, color=GREY_DARK)
    bullets = [
        "SWIR alone outperforms VNIR alone by +5.94 pp OA (90.85 → 96.79).",
        "Marginal gain adding VNIR to SWIR: +0.59 pp.",
        "Implication for this thesis: SuperDove has no SWIR. How much can the 5 extra VNIR bands (vs RGB) recover?",
    ]
    for line in bullets:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=12, bullet=True, color=BLACK)
    add_refs(s, "Aguilar, Jiménez-Lao, Aguilar 2021 — Remote Sensing 13(11), 2133. DOI: 10.3390/rs13112133.")
    add_page_num(s, 12, TOTAL)

    # --- Slide 13 — Where diagnostic info lives ---
    s = new_slide(prs)
    add_title(s, "Where the diagnostic information lives")
    add_image(s, ASSETS / "12_diagnostic_regions.png",
              Inches(0.5), Inches(1.2), width=Inches(12.3))
    add_body(s, [
        ("SuperDove ends at 0.86 µm; the most chemistry-diagnostic regions live in SWIR.", True, 0),
        ("Open hypothesis for this thesis: do extra VNIR bands (Coastal Blue, Yellow, Red Edge, NIR2) recover part of the gap?", True, 0),
    ], top=Inches(6.0), size=11)
    add_refs(s, "Spectra: USGS splib07a (Kokaly 2017). Diagnostic features: Cilia 2015, Zhou 2021, Aguilar 2025.")
    add_page_num(s, 13, TOTAL)

    # --- Slide 14 — More bands ≠ more information ---
    s = new_slide(prs)
    add_title(s, "More bands ≠ more information — well-chosen few suffice")
    add_image(s, ASSETS / "13_bands_vs_oa.png",
              Inches(0.5), Inches(1.3), width=Inches(8.5))
    tb = s.shapes.add_textbox(Inches(9.3), Inches(1.5),
                               Inches(3.8), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0], "Measured evidence",
                  size=12, bold=True, color=BLACK)
    for line in [
        "Vitek 2025 (C&D waste):",
        "  RGB + 2 narrowbands ≈ HSI 768 b. Plateau at 5 bands.",
        "",
        "Aguilar 2021 (WV-3 plastic):",
        "  VNIR (8 b) 90.85 → All (16 b) 97.38.",
        "  SWIR carries most of the jump.",
        "",
        "Zhou 2021 (WV-3 plastic):",
        "  8 SWIR bands discriminate 3 polymer clusters.",
        "",
        "Caveat. Evidence is SWIR-heavy.",
        "VNIR-only is the open question this thesis tests.",
    ]:
        p = tf.add_paragraph()
        is_caveat = "Caveat" in line or "this thesis" in line
        set_paragraph(p, line, size=10,
                      italic=is_caveat,
                      color=GREY_DARK if not is_caveat else BLACK)
    add_refs(s, "Vitek et al. 2025 (CTU Prague) · Aguilar et al. 2021 · Zhou et al. 2021.")
    add_page_num(s, 14, TOTAL)

    # --- Slide 15 — SuperDove the chosen trade-off ---
    s = new_slide(prs)
    add_title(s, "SuperDove: the chosen trade-off for this thesis")
    add_image(s, ASSETS / "14_superdove_bands.png",
              Inches(0.5), Inches(1.2), width=Inches(12.3))
    # 2-column text below
    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.3),
                               Inches(6.0), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0], "What you get",
                  size=12, bold=True, color=BLACK)
    for line in [
        "• 3 m GSD — ~11× more pixels per m² than Sentinel-2",
        "• 8 VNIR bands (Coastal, Blue, 2×Green, Yellow, Red, RE, NIR)",
        "• Sub-daily revisit",
        "• Free for accredited E&R users",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=10.5, color=BLACK)

    tb2 = s.shapes.add_textbox(Inches(6.8), Inches(5.3),
                                Inches(6.0), Inches(1.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2.paragraphs[0], "Known constraints",
                  size=12, bold=True, color=BLACK)
    for line in [
        "• No SWIR — chemistry-diagnostic region missing",
        "• Radiometric stability across constellation varies (~6% mean err. vs S-2)",
        "• 8-band cubes still a young setting for FM transfer",
    ]:
        p = tf2.add_paragraph()
        set_paragraph(p, line, size=10.5, color=BLACK)

    add_refs(s, "Planet PSScene SuperDove specs · radiometry: Wang 2023.")
    add_page_num(s, 15, TOTAL)

    # --- Slide 16 — FM intro ---
    s = new_slide(prs)
    add_title(s, "Foundation models for Earth Observation: a new lever")
    add_body(s, [
        ("2024–2025: a wave of pretrained backbones for EO.", True, 0),
        ("Most are pretrained at 10–30 m on Sentinel-2 / HLS — naïve transfer to 3 m SuperDove is not guaranteed.", True, 0),
        ("Two paths forward:", True, 0),
        ("Sensor-agnostic models (DOFA, AnySat) — designed to accept arbitrary bands.", True, 1),
        ("Parameter-efficient adapters (DEFLECT) — add few-parameter heads on top of any pretrained ViT.", True, 1),
        ("Open question this thesis answers: how well does pretrain at 10–30 m transfer to 3 m?", True, 0),
    ], top=Inches(1.4), size=15)
    add_refs(s, "FM-RS survey 2024 · Xiong 2024 (DOFA) · Astruc 2024 (AnySat) · Thoreau 2025 (DEFLECT) · Szwarcman 2024 (Prithvi-EO-2.0).")
    add_page_num(s, 16, TOTAL)

    # --- Slide 17 — DOFA dedicated ---
    s = new_slide(prs)
    add_title(s, "DOFA: a band-agnostic backbone — the candidate to start from")
    # Left: text
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5),
                               Inches(7.5), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0],
                  "Dynamic One-For-All (Xiong et al. 2024)",
                  size=14, bold=True, color=BLACK)
    p2 = tf.add_paragraph()
    set_paragraph(p2, "What it does",
                  size=12, bold=True, color=BLACK)
    for line in [
        "A hypernetwork generates patch-embedding weights conditioned on each input band's central wavelength.",
        "One backbone handles S-2 (13 b), SuperDove (8 b), WV-3 (16 b), or HSI.",
        "Generalises to unseen sensor configurations at inference.",
        "Plug-and-play for SuperDove via the 8 central λ; no retraining needed.",
        "Natural fit for the band-ablation study (RGB / +NIR / +RedEdge / 8 b).",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=11, bullet=True, color=BLACK)

    pq = tf.add_paragraph()
    set_paragraph(pq, "What to verify",
                  size=12, bold=True, color=BLACK)
    for line in [
        "Does pretrain at 10–30 m close the gap to 3 m? (open, empirical)",
        "Combinable with DEFLECT adapter (<1 % params) for spectral fine-tune.",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=11, bullet=True, color=BLACK)

    # Right: DOFA paper architecture figure (cropped)
    add_image(s, PDF_FIGS / "16_dofa_fig2_crop.png",
              Inches(8.5), Inches(1.4), height=Inches(5.0))

    add_refs(s, "Xiong et al. 2024, arXiv:2403.15356 (DOFA). Adapter: Thoreau et al. 2025 (DEFLECT, ICCV).")
    add_page_num(s, 17, TOTAL)

    # --- Slide 18 — Shepherd asbestos competitor ---
    s = new_slide(prs)
    add_title(s, "The closest competitor: spaceborne hyperspectral asbestos detection")
    # Left text
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5),
                               Inches(6.8), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf.paragraphs[0],
                  "Shepherd et al. 2025 — Scientific Reports (Nature OA)",
                  size=13, bold=True, color=BLACK)
    p2 = tf.add_paragraph()
    set_paragraph(p2, "What they did",
                  size=12, bold=True, color=BLACK)
    for line in [
        "EnMAP Level 2A hyperspectral (230 bands, VNIR+SWIR, 30 m).",
        "Cascade of 8 supervised classifiers (LSU, SVM, SAM, ACE, MD, ML, SID, MF).",
        "Field-spectra calibration with ASD FieldSpec 4 (ground-truth survey).",
        "Result: 86 % positive-match rate; ACE 91.4 %, SID 90.1 %, SVM 89.2 % OA.",
        "Works because EnMAP has SWIR: Mg-OH chrysotile feature @ 2.31 µm.",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=11, bullet=True, color=BLACK)

    p3 = tf.add_paragraph()
    set_paragraph(p3, "How this thesis differs",
                  size=12, bold=True, color=BLACK)
    for line in [
        "SuperDove VNIR 8 b @ 3 m vs EnMAP HSI 230 b @ 30 m.",
        "10× finer spatially, no SWIR — opposite trade-off.",
        "Asks: can VNIR signal + spatial resolution compete?",
    ]:
        p = tf.add_paragraph()
        set_paragraph(p, line, size=11, bullet=True, color=BLACK)

    add_image(s, PDF_FIGS / "17_shepherd_fig2_crop.png",
              Inches(7.5), Inches(1.4), height=Inches(5.0))

    add_refs(s, "Shepherd, Sagi, Zagron et al. 2025 — Sci. Rep. 15:24166. DOI: 10.1038/s41598-025-09738-w.")
    add_page_num(s, 18, TOTAL)

    # --- Slide 19 — Gaps + the test that matters ---
    s = new_slide(prs)
    add_title(s, "Gaps in the literature — and the test that matters")
    add_body(s, [
        ("Open gaps", False, 0),
        ("Material-level labels are still missing — no public terrestrial multispectral waste dataset annotates the material.", True, 0),
        ("RGB pipelines conflate iso-colour, chemically distinct materials.", True, 0),
        ("No cross-sensor / cross-region benchmark for waste materials.", True, 0),
        ("Generalisation is tested anecdotally, not systematically.", True, 0),
        ("No public SuperDove-era (post-2020) waste benchmark.", True, 0),
        ("EO foundation models pretrained at 10–30 m: transfer to 3 m unproven.", True, 0),
    ], top=Inches(1.4), size=13)

    # Quote box
    qb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(5.6),
                             Inches(12.3), Inches(1.2))
    qb.fill.background()
    qb.line.color.rgb = BLACK
    qb.line.width = Pt(1.2)
    qtf = qb.text_frame
    qtf.margin_left = qtf.margin_right = Inches(0.3)
    qtf.margin_top = qtf.margin_bottom = Inches(0.15)
    qtf.word_wrap = True
    set_paragraph(qtf.paragraphs[0],
                  "The test that matters",
                  size=12, bold=True, color=BLACK)
    pq = qtf.add_paragraph()
    set_paragraph(pq,
                  "Does material classification improve generalisation across sites?",
                  size=15, italic=True, color=BLACK)
    add_refs(s, "Gap list distilled from Fraternali et al. 2024 (PoliMi survey, arXiv:2402.09066).")
    add_page_num(s, 19, TOTAL)

    # --- Slide 20 — Proposed direction (roadmap) ---
    s = new_slide(prs)
    add_title(s, "Proposed direction: from controlled pilot to public benchmark")
    add_image(s, ASSETS / "19_roadmap.png",
              Inches(0.7), Inches(1.4), width=Inches(11.9))
    add_body(s, [
        ("Phase 1 (asbestos pilot) tests the question first on a known-difficult material with public WFS ground truth.", True, 0),
        ("Phase 2 (MS waste benchmark) generalises the answer to multi-class material classification.", True, 0),
    ], top=Inches(5.5), size=12)
    add_refs(s, "Pilot: Mappatura 2020 (Lombardia WFS). Backbone: Xiong 2024 (DOFA). Baseline: Gibellini 2025.")
    add_page_num(s, 20, TOTAL)

    prs.save(OUT)
    print(f"✓ Deck saved: {OUT}")
    return OUT


if __name__ == "__main__":
    build()
