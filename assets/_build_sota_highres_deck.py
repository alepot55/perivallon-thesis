#!/usr/bin/env python3
"""Build the focused high-resolution material-discrimination SOTA deck.

Thomas / internal register: absolute B/N, no color, solid bullets, English,
self-explanatory titles, per-work table rows, refs as title + year.
Source content: docs/02_research/sota_highres_material.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x66, 0x66, 0x66)
LINE = RGBColor(0xBB, 0xBB, 0xBB)

FONT = "Arial"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

M = Inches(0.7)  # margin


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def title(slide, text, sub=None):
    tf = _tb(slide, M, Inches(0.5), EMU_W - 2 * M, Inches(1.0))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BLACK
    # underline rule
    ln = slide.shapes.add_connector(2, M, Inches(1.45), EMU_W - M, Inches(1.45))
    ln.line.color.rgb = BLACK; ln.line.width = Pt(1.5)
    if sub:
        tf2 = _tb(slide, M, Inches(1.55), EMU_W - 2 * M, Inches(0.6))
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = sub
        r2.font.name = FONT; r2.font.size = Pt(15); r2.font.italic = True; r2.font.color.rgb = GREY


def bullets(slide, items, top=Inches(1.9), size=18, gap=10):
    tf = _tb(slide, M, top, EMU_W - 2 * M, EMU_H - top - Inches(0.5))
    first = True
    for it in items:
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            lvl, txt = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        # solid bullet, never dashes
        bullet = "•  " if lvl == 0 else "◦  "
        r = p.add_run(); r.text = bullet + txt
        r.font.name = FONT
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = BLACK if lvl == 0 else RGBColor(0x33, 0x33, 0x33)


def footer(slide, text):
    tf = _tb(slide, M, EMU_H - Inches(0.45), EMU_W - 2 * M, Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


def table_slide(slide, headers, rows, col_w, top=Inches(1.9), fsize=12):
    nrows = len(rows) + 1
    ncols = len(headers)
    total_w = sum(col_w)
    left = M
    height = Inches(0.0)
    gtbl = slide.shapes.add_table(nrows, ncols, left, top, Emu(int(total_w)), Inches(0.4 * nrows))
    tbl = gtbl.table
    # disable banding/style colors -> we set manually
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = Emu(int(w))
    # header
    for ci, htxt in enumerate(headers):
        c = tbl.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Pt(4); c.margin_right = Pt(4); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = htxt
        r.font.name = FONT; r.font.size = Pt(fsize); r.font.bold = True; r.font.color.rgb = BLACK
    # body
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Pt(4); c.margin_right = Pt(4); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT; r.font.size = Pt(fsize); r.font.color.rgb = BLACK
            if ci == 0:
                r.font.bold = True
    return tbl


def new():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    return s


# ---- Slide 1: title ----
s = new()
tf = _tb(s, M, Inches(2.4), EMU_W - 2 * M, Inches(2.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "High-resolution material discrimination from remote sensing"
r.font.name = FONT; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = BLACK
p2 = tf.add_paragraph(); p2.space_before = Pt(14)
r2 = p2.add_run(); r2.text = "A focused state of the art for the SuperDove (8 VNIR bands, 3 m) setting"
r2.font.name = FONT; r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
ln = s.shapes.add_connector(2, M, Inches(4.55), Inches(6.5), Inches(4.55))
ln.line.color.rgb = BLACK; ln.line.width = Pt(1.5)
footer(s, "PERIVALLON thesis · re-scoped after the 2026-06-15 advisor directive")

# ---- Slide 2: the question, re-scoped ----
s = new()
title(s, "The question, re-scoped", "Not “waste from satellites” — how the literature discriminates MATERIALS at a useful resolution")
bullets(s, [
    "Input we operate on: SuperDove — 8 VNIR bands, 3 m GSD, no SWIR.",
    "Task we care about: identify what a surface is made of (material-level), per-pixel / per-object / per-region.",
    "Works on other materials count — minerals, roofing, urban surfaces — because the question is material discrimination, not waste detection.",
    "Why this matters: it places the thesis against the right comparators, instead of a broad and unfocused “waste RS” survey.",
])

# ---- Slide 3: inclusion / exclusion ----
s = new()
title(s, "What counts — and what we deliberately exclude")
bullets(s, [
    "A work is included only if it satisfies all three:",
    (1, "Task = material discrimination (not scene classification, not binary waste/no-waste, not land-use)."),
    (1, "GSD ≤ ~5 m (sub-metre VHR, and the ~3 m PlanetScope / SuperDove class)."),
    (1, "Multispectral or hyperspectral input (≥ 4 bands)."),
    "Excluded by resolution — even with SWIR or material labels:",
    (1, "Sentinel-2 (10 m), EnMAP / PRISMA / EMIT (30–60 m)."),
    (1, "Known, not missed: Shepherd 2025 (EnMAP asbestos), Estrela 2025 (EMIT plastic), MARIDA / MADOS — they set the SWIR upper bound at a resolution we ruled out.",),
], size=17)

# ---- Slide 4: the core tension ----
s = new()
title(s, "The core tension: SWIR discriminates, resolution locates")
bullets(s, [
    "Material identity is carried by SWIR features (e.g. asbestos Mg-OH at 2.31 µm, polymer C–H at ~1215 / 1730 nm).",
    "Spatial resolution decides whether you can place those materials on the ground.",
    "On satellites, high resolution + SWIR comes either expensive and access-limited (WorldView-3, 3.7 m SWIR) or at low resolution (EnMAP / PRISMA, 30 m).",
    "SuperDove (8 VNIR @ 3 m, no SWIR) sits in a sparsely-studied corner — between sub-metre VHR and 10 m Sentinel-2.",
    "Cleanest single number for what SWIR buys: Saba 2025, same asbestos task — HS 97.3% vs MS 74.4%.",
])
footer(s, "Saba et al. 2025 (RSASE) · same asbestos task, hyperspectral vs multispectral")

# ---- Slide 5: our operating point = empty cell ----
s = new()
title(s, "Where everyone sits — and the cell nobody fills")
bullets(s, [
    "Across 24 verified works, only three are at once material-discrimination, ≤ 3 m, and VNIR-only like SuperDove:",
    (1, "Saba 2026 — WV-3 VNIR, 1.24 m."),
    (1, "Widipaminto 2021 — Pléiades, 2 m."),
    (1, "Xu 2022 — UAV, cm-scale."),
    "All three are finer than 3 m. Every strong polymer- or mineral-identity result has SWIR.",
    "No work sits at our exact operating point: 3 m, 8-band VNIR, no SWIR.",
    "That empty cell is the thesis.",
])

# ---- Slide 6: Domain 1 plastics ----
s = new()
title(s, "Plastics & polymers")
cw = [Inches(3.3), Inches(2.7), Inches(2.4), Inches(3.5)]
rows = [
    ["Aguilar 2025", "WV-3 / 3.7 m · SWIR ✓", "matched-filter, 5 polymers", "precision 92.5%"],
    ["Zhou 2021", "WV-3 / 3.7 m · SWIR ✓", "pixel class., 3 polymers", "OA > 80%"],
    ["Alboody 2023", "aquatic drone / cm · SWIR ✓", "supervised, up to 10 polymers", "OA ≈ 89%"],
    ["Agronomy 2026", "UAV / cm · VNIR ✗", "seg: mulch / soil / cotton", "mIoU 85.9% > RGB 83.4%"],
]
table_slide(s, ["Work (year)", "Platform / GSD · SWIR", "Task / method", "Headline metric"], rows, cw)
tf = _tb(s, M, Inches(4.6), EMU_W - 2 * M, Inches(2.2))
bl = ["Takeaway: polymer-TYPE identity needs bands beyond 1000 nm (SWIR).",
      "But VNIR alone still separates plastic presence/extent and beats RGB (Agronomy 2026).",
      "For SuperDove: plausible for plastic presence, not for polymer identity."]
for i, t in enumerate(bl):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = "•  " + t
    r.font.name = FONT; r.font.size = Pt(15); r.font.color.rgb = BLACK

# ---- Slide 7: Domain 2 asbestos ----
s = new()
title(s, "Asbestos-cement & roofing — closest to the Phase-1 pilot")
rows = [
    ["Bonifazi 2026", "WV-3 / 1.24 m + SWIR ✓", "Py6S + MLC, building-level", "F1 0.87"],
    ["Cartagena 2024", "airborne HySpex / 0.8 m · SWIR ✓", "SAM, AC vs clay vs steel", "OA 96%, AC 98%"],
    ["Saba 2026", "WV-3 / 1.24 m · VNIR ✗", "32 classifiers, multiclass", "Macro-F1 97.6%"],
    ["Saba 2025", "WV-3 MS + airborne HS · ✗/✓", "MS vs HS comparison", "HS 97.3% vs MS 74.4%"],
    ["Widipaminto 2021", "Pléiades / 2 m · VNIR ✗", "SVM, 5 roof materials", "OA 92.9%, κ 0.91"],
    ["Raczko 2022", "airborne ortho / 25 cm · VNIR ✗", "CNN, AC binary", "OA 88–93%"],
]
table_slide(s, ["Work (year)", "Platform / GSD · SWIR", "Task / method", "Headline metric"], rows, cw, top=Inches(1.7), fsize=12)
tf = _tb(s, M, Inches(5.3), EMU_W - 2 * M, Inches(1.7))
bl = ["The VNIR/RGB line is now strong and quantified — best external evidence the VNIR-only pilot is viable.",
      "Two caveats: all operate at 0.25–2 m (12–25× finer than 3 m); Saba 2025 quantifies the VNIR penalty (~23 pts)."]
for i, t in enumerate(bl):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = "•  " + t
    r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = BLACK

# ---- Slide 8: Domain 3 inert/minerals ----
s = new()
title(s, "Inert / C&D waste & minerals")
rows = [
    ["Xu 2022", "UAV / cm · VNIR ✗", "decision tree, 4 C&D classes", "OA 85.9%, κ 0.85"],
    ["Vitek/Krauz 2025", "close-range / lab · RGB+narrow", "band selection for C&D", "RGB + 2 bands ≈ 768-band HSI"],
    ["Karimzadeh 2021", "WV-3 / 1.2 m + SWIR ✓", "SVM, lithologies", "OA 88.4% (top band = SWIR-7)"],
    ["Wang 2022", "airborne HySpex / 1 m · SWIR ✓", "RF, 7 alteration minerals", "OA 73.1%"],
]
table_slide(s, ["Work (year)", "Platform / GSD · SWIR", "Task / method", "Headline metric"], rows, cw)
tf = _tb(s, M, Inches(4.6), EMU_W - 2 * M, Inches(2.2))
bl = ["Xu 2022 is the closest sensor-analogue: VNIR-only inert/C&D, OA 85.9% with a simple tree.",
      "“Which bands, not how many” (Vitek/Krauz) is the strongest enabler for 8 VNIR bands.",
      "The Chinese mineral cluster (Wang / Huang / Huaniushan) is real — but airborne HSI + SWIR, none at 3 m VNIR-only, none on waste/asbestos. The gap survives the advisor’s check."]
for i, t in enumerate(bl):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = "•  " + t
    r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = BLACK

# ---- Slide 9: Domain 4 metals/urban ----
s = new()
title(s, "Metals & urban materials — confusers & baselines")
rows = [
    ["Ma 2023", "airborne CASI / 2.5 m · VNIR ✗", "3D-CNN, 15 urban classes", "OA 96.3%"],
    ["Banolia 2022", "airborne HSI / 1.3–2 m · VNIR ✗", "rule-based metal-roof detection", "OA 99% (F1 0.44, GT noise)"],
    ["MCubeS", "close-range / cm · RGB+NIR+pol", "20-class material seg", "NIR+pol adds +2.7–8.7% mIoU"],
    ["Toulouse 2024", "airborne / 1 m · SWIR ✓", "32 urban-material classes", "public benchmark"],
]
table_slide(s, ["Work (year)", "Platform / GSD · SWIR", "Task / method", "Headline metric"], rows, cw)
tf = _tb(s, M, Inches(4.6), EMU_W - 2 * M, Inches(2.0))
bl = ["Metal’s high-reflectance, flat-spectrum cue survives in VNIR (Banolia) — metal sheeting is a common waste/roof confuser.",
      "Ma 2023: a VNIR 3D-CNN reaches OA 96% at 2.5 m — a direct methodological analogue for a 3 m SuperDove CNN baseline."]
for i, t in enumerate(bl):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    r = p.add_run(); r.text = "•  " + t
    r.font.name = FONT; r.font.size = Pt(14); r.font.color.rgb = BLACK

# ---- Slide 10: foundation models ----
s = new()
title(s, "Foundation models at high resolution", "The honest assessment")
bullets(s, [
    "DOFA (Xiong 2024, preprint) conditions on wavelength only, not on GSD.",
    (1, "Closes the band-mismatch gap — SuperDove’s 8 VNIR centres are accepted natively."),
    (1, "Does not close the GSD / object-scale gap — that stays an open, measurable question."),
    "PANGAEA (Marsocci 2024): geo-FMs do not consistently beat supervised baselines, incl. high-res — the strongest skeptic-supporting citation.",
    "Testable mitigations: Scale-MAE (GSD-aware), DEFLECT (RGB→MS PEFT), two-branch RGB+MS fusion.",
    "Verdict: treat FM adaptation as a bounded side test; the band-ablation on a tuned supervised baseline is the defensible core.",
], top=Inches(2.3), size=17)

# ---- Slide 11: synthesis ----
s = new()
title(s, "Cross-cutting synthesis")
bullets(s, [
    "SWIR is the discriminator; resolution is the locator. Cleanest number: Saba 2025, HS 97.3% vs MS 74.4%.",
    "VNIR-only at high resolution genuinely works — asbestos, inert/C&D, urban, plastic-presence. It fails for polymer identity.",
    "It is “which bands, not how many” (Vitek/Krauz: RGB + 2 bands ≈ 768-band HSI).",
    "The catch is GSD, not bands: almost every VNIR success is at 0.25–2.5 m, none at 3 m.",
    "The Chinese literature does not pre-empt us — it is airborne HSI + SWIR mineral mapping.",
    "3 m VNIR-only is the frontier, not the comfort zone — the thesis’s contribution space.",
], size=17)

# ---- Slide 12: where the thesis sits ----
s = new()
title(s, "Where the thesis sits & the test that matters")
bullets(s, [
    "Setting: SuperDove, 8 VNIR bands, 3 m, no SWIR — between sub-metre VHR and 10 m Sentinel-2, essentially unstudied.",
    "The gap: no high-res material-discrimination study operates at exactly VNIR-only 3 m; the closest precedents are all finer.",
    "The test that matters: how much of the SWIR/VHR material-discrimination performance can VNIR-only 3 m recover?",
    (1, "Quantified by the band ablation (RGB → +NIR → +RedEdge → 8 bands), benchmarked against the VNIR competitors above."),
    (1, "Most honest yardstick for what SWIR buys: Saba 2025, 97.3% (HS) vs 74.4% (MS) on the asbestos task."),
], size=18)

prs.save("/home/alepot55/Desktop/uni/Tesi/sota_highres_material_deck.pptx")
print("saved sota_highres_material_deck.pptx —", len(prs.slides._sldIdLst), "slides")
